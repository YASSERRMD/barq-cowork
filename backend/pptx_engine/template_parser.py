"""
Template parser: extracts design information from an existing PPTX file.
Reads theme XML, font schemes, layout placeholders, and auto-generates brand rules.
"""

from __future__ import annotations

import re
from typing import Optional

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from .schema import (
    BrandRule,
    LayoutInfo,
    MasterInfo,
    ParsedTemplate,
    PlaceholderInfo,
    ThemeColors,
    ThemeFonts,
    ValidationSeverity,
)


_HEX_RE = re.compile(r"^[0-9A-Fa-f]{6}$")

# Known theme color slot names → schema field names
_COLOR_SLOT_MAP = {
    "dk1": "text",
    "lt1": "background",
    "dk2": "surface",
    "lt2": "text_muted",
    "accent1": "accent",
    "accent2": "accent2",
    "accent3": "border",
    "accent4": "success",
    "accent5": "warning",
    "accent6": "danger",
}


def _emu_to_int(val) -> int:
    """Safely convert EMU-like value to plain int."""
    try:
        return int(val)
    except (TypeError, ValueError):
        return 0


class TemplateParser:
    """Parse an existing PPTX file and extract all design information."""

    def parse(self, pptx_path: str) -> ParsedTemplate:
        """
        Load a PPTX file and extract theme, fonts, layouts, and master info.

        Returns a fully populated ParsedTemplate.
        """
        prs = Presentation(pptx_path)

        theme_colors = self._extract_theme_colors(prs)
        theme_fonts = self._extract_theme_fonts(prs)
        layouts = self._extract_layouts(prs)
        masters = self._extract_masters(prs)

        slide_width = _emu_to_int(prs.slide_width)
        slide_height = _emu_to_int(prs.slide_height)

        return ParsedTemplate(
            theme_colors=theme_colors,
            theme_fonts=theme_fonts,
            layouts=layouts,
            masters=masters,
            slide_count=len(prs.slides),
            slide_width=slide_width,
            slide_height=slide_height,
            source_path=pptx_path,
        )

    def _extract_theme_colors(self, prs: Presentation) -> ThemeColors:
        """
        Extract theme colors from the first slide master's theme XML (a:clrScheme).
        Falls back to defaults when elements are missing.
        """
        color_map: dict[str, str] = {}

        try:
            for master in prs.slide_masters:
                theme_el = master.element.find(
                    ".//" + qn("a:theme")
                )
                if theme_el is None:
                    # Try theme through theme part relationship
                    try:
                        theme_el = master.theme_color_map._element  # type: ignore[attr-defined]
                    except AttributeError:
                        pass

                # Walk the XML tree looking for a:clrScheme
                root = master.element
                clr_scheme = root.find(".//" + qn("a:clrScheme"))
                if clr_scheme is None:
                    continue

                for child in clr_scheme:
                    tag_local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                    slot_name = _COLOR_SLOT_MAP.get(tag_local)
                    if slot_name is None:
                        continue

                    hex_val = self._color_element_to_hex(child)
                    if hex_val:
                        color_map[slot_name] = f"#{hex_val.upper()}"

                if color_map:
                    break

        except Exception:
            pass  # Fall through to defaults

        # Build ThemeColors with extracted values merged over defaults
        defaults = ThemeColors()
        kwargs = defaults.model_dump()
        kwargs.update(color_map)
        return ThemeColors(**kwargs)

    def _color_element_to_hex(self, el) -> Optional[str]:
        """
        Resolve a color element (a:srgbClr, a:sysClr, a:dk1, etc.) to a hex string.
        Returns 6-char uppercase hex or None.
        """
        # Direct srgbClr child
        srgb = el.find(qn("a:srgbClr"))
        if srgb is not None:
            val = srgb.get("val", "")
            if _HEX_RE.match(val):
                return val.upper()

        # sysClr has lastClr attribute
        sys = el.find(qn("a:sysClr"))
        if sys is not None:
            last = sys.get("lastClr", "")
            if _HEX_RE.match(last):
                return last.upper()

        # schemeClr — not resolvable without full theme, skip
        return None

    def _extract_theme_fonts(self, prs: Presentation) -> ThemeFonts:
        """
        Extract font family names from a:fontScheme inside the first slide master.
        """
        heading_font = "Calibri Light"
        body_font = "Calibri"

        try:
            for master in prs.slide_masters:
                root = master.element
                font_scheme = root.find(".//" + qn("a:fontScheme"))
                if font_scheme is None:
                    continue

                major = font_scheme.find(qn("a:majorFont"))
                if major is not None:
                    latin = major.find(qn("a:latin"))
                    if latin is not None:
                        typeface = latin.get("typeface", "")
                        if typeface and typeface not in ("+mj-lt", "+mj-cs"):
                            heading_font = typeface

                minor = font_scheme.find(qn("a:minorFont"))
                if minor is not None:
                    latin = minor.find(qn("a:latin"))
                    if latin is not None:
                        typeface = latin.get("typeface", "")
                        if typeface and typeface not in ("+mn-lt", "+mn-cs"):
                            body_font = typeface
                break

        except Exception:
            pass

        return ThemeFonts(
            heading=heading_font,
            body=body_font,
        )

    def _extract_layouts(self, prs: Presentation) -> list[LayoutInfo]:
        """
        Extract all slide layouts from all masters with placeholder info.
        """
        layouts: list[LayoutInfo] = []
        seen_names: set[str] = set()

        slide_width = _emu_to_int(prs.slide_width)
        slide_height = _emu_to_int(prs.slide_height)

        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                name = layout.name or "Unknown"
                # Deduplicate by name
                unique_name = name
                idx = 1
                while unique_name in seen_names:
                    unique_name = f"{name}_{idx}"
                    idx += 1
                seen_names.add(unique_name)

                placeholders: list[PlaceholderInfo] = []
                for ph in layout.placeholders:
                    try:
                        ph_info = PlaceholderInfo(
                            idx=ph.placeholder_format.idx,
                            placeholder_type=str(ph.placeholder_format.type),
                            left=_emu_to_int(ph.left),
                            top=_emu_to_int(ph.top),
                            width=_emu_to_int(ph.width),
                            height=_emu_to_int(ph.height),
                            name=ph.name,
                        )
                        placeholders.append(ph_info)
                    except Exception:
                        continue

                layouts.append(
                    LayoutInfo(
                        name=unique_name,
                        placeholders=placeholders,
                        slide_width=slide_width,
                        slide_height=slide_height,
                    )
                )

        return layouts

    def _extract_masters(self, prs: Presentation) -> list[MasterInfo]:
        """Extract slide master metadata."""
        masters: list[MasterInfo] = []
        for i, master in enumerate(prs.slide_masters):
            name = f"Master {i + 1}"
            try:
                # Try to get master name from XML
                nm_el = master.element.find(".//" + qn("p:cSld"))
                if nm_el is not None:
                    nm = nm_el.get("name", "")
                    if nm:
                        name = nm
            except Exception:
                pass

            masters.append(
                MasterInfo(
                    name=name,
                    layout_count=len(master.slide_layouts),
                )
            )

        return masters

    def generate_brand_rules(self, parsed: ParsedTemplate) -> list[BrandRule]:
        """
        Auto-generate brand compliance rules from a parsed template.
        Rules cover font families, background colors, and accent usage.
        """
        rules: list[BrandRule] = []

        # Font rules
        if parsed.theme_fonts.heading:
            rules.append(
                BrandRule(
                    name="Heading Font",
                    rule_type="font",
                    condition="Slide headings must use the template heading font",
                    allowed_values=[parsed.theme_fonts.heading, "Calibri Light", "Calibri"],
                    severity=ValidationSeverity.warning,
                )
            )

        if parsed.theme_fonts.body:
            rules.append(
                BrandRule(
                    name="Body Font",
                    rule_type="font",
                    condition="Body text must use the template body font",
                    allowed_values=[parsed.theme_fonts.body, "Calibri", "Calibri Light"],
                    severity=ValidationSeverity.warning,
                )
            )

        # Background color rule
        bg_color = parsed.theme_colors.background
        if bg_color:
            rules.append(
                BrandRule(
                    name="Background Color",
                    rule_type="color",
                    condition="Slide background must match template background color",
                    allowed_values=[bg_color, "#0F172A", "#FFFFFF", "#F8FAFC"],
                    severity=ValidationSeverity.info,
                )
            )

        # Accent color rule
        accent_color = parsed.theme_colors.accent
        if accent_color:
            rules.append(
                BrandRule(
                    name="Primary Accent Color",
                    rule_type="color",
                    condition="Primary accent elements must use the template accent color",
                    allowed_values=[accent_color, parsed.theme_colors.accent2],
                    severity=ValidationSeverity.warning,
                )
            )

        # Slide count guideline
        rules.append(
            BrandRule(
                name="Slide Count",
                rule_type="layout",
                condition="Presentations should have between 5 and 30 slides",
                allowed_values=["5-30"],
                severity=ValidationSeverity.info,
            )
        )

        # Content rules
        rules.append(
            BrandRule(
                name="Bullet Point Limit",
                rule_type="content",
                condition="Bullet slides should not exceed 6 points",
                allowed_values=["1", "2", "3", "4", "5", "6"],
                severity=ValidationSeverity.warning,
            )
        )

        rules.append(
            BrandRule(
                name="Heading Length",
                rule_type="content",
                condition="Slide headings should not exceed 60 characters",
                allowed_values=["max:60"],
                severity=ValidationSeverity.warning,
            )
        )

        return rules
