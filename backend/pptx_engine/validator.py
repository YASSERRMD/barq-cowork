"""
Deck validation: schema, content-length, rendered overflow, and brand compliance.
"""

from __future__ import annotations

import io
import re
from typing import Optional

from pptx import Presentation
from pptx.util import Pt

from .schema import (
    BrandRule,
    Deck,
    SlideType,
    ValidationIssue,
    ValidationResult,
    ValidationSeverity,
)


class DeckValidator:
    """Multi-layer deck validator."""

    def __init__(self, brand_rules: Optional[list[BrandRule]] = None):
        self._brand_rules = brand_rules or []

    # ── Public entry point ─────────────────────────────────────────────────────

    def validate(self, deck: Deck, pptx_bytes: Optional[bytes] = None) -> ValidationResult:
        issues: list[ValidationIssue] = []
        issues += self._validate_schema(deck)
        issues += self._validate_slide_id_uniqueness(deck)
        issues += self._validate_content_length(deck)
        if pptx_bytes:
            issues += self._validate_rendered_overflow(pptx_bytes)
            issues += self._validate_brand_compliance(deck, pptx_bytes)
        return ValidationResult(issues=issues)

    # ── Schema validation ──────────────────────────────────────────────────────

    def _validate_schema(self, deck: Deck) -> list[ValidationIssue]:
        issues: list[ValidationIssue] = []

        if not deck.meta.title or not deck.meta.title.strip():
            issues.append(ValidationIssue(
                severity=ValidationSeverity.error,
                code="MISSING_TITLE",
                message="Deck must have a non-empty title.",
            ))

        if not deck.plan.slides:
            issues.append(ValidationIssue(
                severity=ValidationSeverity.error,
                code="NO_SLIDES",
                message="Deck must contain at least one slide.",
            ))
            return issues

        if len(deck.plan.slides) > 50:
            issues.append(ValidationIssue(
                severity=ValidationSeverity.warning,
                code="TOO_MANY_SLIDES",
                message=f"Deck has {len(deck.plan.slides)} slides; consider splitting at 50.",
            ))

        for slide in deck.plan.slides:
            if not slide.id:
                issues.append(ValidationIssue(
                    severity=ValidationSeverity.error,
                    code="MISSING_SLIDE_ID",
                    message=f"Slide is missing an ID (heading: '{slide.heading[:30]}').",
                    slide_id=slide.id,
                ))

            if len(slide.heading) > 60:
                issues.append(ValidationIssue(
                    severity=ValidationSeverity.warning,
                    code="HEADING_TOO_LONG",
                    message=(
                        f"Slide heading is {len(slide.heading)} chars; max recommended is 60."
                    ),
                    slide_id=slide.id,
                    field="heading",
                ))

            if slide.type == SlideType.title and not slide.heading:
                issues.append(ValidationIssue(
                    severity=ValidationSeverity.warning,
                    code="TITLE_SLIDE_NO_HEADING",
                    message="Title slide should have a heading.",
                    slide_id=slide.id,
                    field="heading",
                ))

        return issues

    # ── Slide ID uniqueness ────────────────────────────────────────────────────

    def _validate_slide_id_uniqueness(self, deck: Deck) -> list[ValidationIssue]:
        issues: list[ValidationIssue] = []
        seen: dict[str, int] = {}
        for i, slide in enumerate(deck.plan.slides):
            if slide.id in seen:
                issues.append(ValidationIssue(
                    severity=ValidationSeverity.error,
                    code="DUPLICATE_SLIDE_ID",
                    message=(
                        f"Slide ID '{slide.id}' is duplicated at positions "
                        f"{seen[slide.id]} and {i}."
                    ),
                    slide_id=slide.id,
                ))
            else:
                seen[slide.id] = i
        return issues

    # ── Content length validation ──────────────────────────────────────────────

    def _validate_content_length(self, deck: Deck) -> list[ValidationIssue]:
        issues: list[ValidationIssue] = []

        for slide in deck.plan.slides:
            c = slide.content

            # Bullet limit
            if c.points and len(c.points) > 6:
                issues.append(ValidationIssue(
                    severity=ValidationSeverity.warning,
                    code="TOO_MANY_BULLETS",
                    message=(
                        f"Slide '{slide.id}' has {len(c.points)} bullet points; max is 6."
                    ),
                    slide_id=slide.id,
                    field="content.points",
                ))

            # Individual bullet length
            if c.points:
                for j, pt in enumerate(c.points):
                    if len(pt) > 200:
                        issues.append(ValidationIssue(
                            severity=ValidationSeverity.info,
                            code="LONG_BULLET",
                            message=(
                                f"Bullet {j+1} on slide '{slide.id}' is {len(pt)} chars; "
                                "consider shortening."
                            ),
                            slide_id=slide.id,
                            field=f"content.points[{j}]",
                        ))

            # Stats count
            if c.stats is not None:
                n = len(c.stats)
                if n < 2:
                    issues.append(ValidationIssue(
                        severity=ValidationSeverity.warning,
                        code="TOO_FEW_STATS",
                        message=f"Stats slide '{slide.id}' has {n} stat(s); recommend 2–4.",
                        slide_id=slide.id,
                        field="content.stats",
                    ))
                if n > 4:
                    issues.append(ValidationIssue(
                        severity=ValidationSeverity.warning,
                        code="TOO_MANY_STATS",
                        message=f"Stats slide '{slide.id}' has {n} stats; max is 4.",
                        slide_id=slide.id,
                        field="content.stats",
                    ))

            # Table column count
            if c.table and c.table.headers:
                if len(c.table.headers) > 8:
                    issues.append(ValidationIssue(
                        severity=ValidationSeverity.warning,
                        code="TOO_MANY_TABLE_COLS",
                        message=(
                            f"Table on slide '{slide.id}' has {len(c.table.headers)} columns; "
                            "max recommended is 8."
                        ),
                        slide_id=slide.id,
                        field="content.table.headers",
                    ))

            # Cards count
            if c.cards and len(c.cards) > 6:
                issues.append(ValidationIssue(
                    severity=ValidationSeverity.warning,
                    code="TOO_MANY_CARDS",
                    message=(
                        f"Cards slide '{slide.id}' has {len(c.cards)} cards; max is 6."
                    ),
                    slide_id=slide.id,
                    field="content.cards",
                ))

            # Steps count
            if c.steps and len(c.steps) > 6:
                issues.append(ValidationIssue(
                    severity=ValidationSeverity.warning,
                    code="TOO_MANY_STEPS",
                    message=(
                        f"Steps slide '{slide.id}' has {len(c.steps)} steps; max is 6."
                    ),
                    slide_id=slide.id,
                    field="content.steps",
                ))

        return issues

    # ── Rendered overflow check ────────────────────────────────────────────────

    def _validate_rendered_overflow(self, pptx_bytes: bytes) -> list[ValidationIssue]:
        """
        Heuristic overflow check: estimate if text volume exceeds shape area.

        Rule: chars_per_point ≈ 10; if text_length * 10pt > shape area (in pt²), warn.
        Shape area is width_pt * height_pt.
        """
        issues: list[ValidationIssue] = []

        try:
            prs = Presentation(io.BytesIO(pptx_bytes))
        except Exception:
            return issues

        OVERFLOW_RATIO = 0.8  # flag when text_area / shape_area > this

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                # Total character count in shape
                total_chars = sum(
                    len(para.text)
                    for para in shape.text_frame.paragraphs
                )
                if total_chars < 50:
                    continue  # short text – no risk

                # Shape area in pt²
                try:
                    w_pt = shape.width / 12700  # EMU → pt
                    h_pt = shape.height / 12700
                except Exception:
                    continue

                if w_pt <= 0 or h_pt <= 0:
                    continue

                shape_area_pt2 = w_pt * h_pt

                # Estimate text area: assume avg char is ~7pt wide × 14pt tall
                avg_char_w = 7
                avg_line_h = 14
                chars_per_line = max(1, int(w_pt / avg_char_w))
                line_count = (total_chars + chars_per_line - 1) // chars_per_line
                text_area_pt2 = line_count * avg_line_h * w_pt

                if text_area_pt2 > shape_area_pt2 * (1 + OVERFLOW_RATIO):
                    slide_note = f"slide[{slide_idx + 1}]"
                    issues.append(ValidationIssue(
                        severity=ValidationSeverity.warning,
                        code="POSSIBLE_OVERFLOW",
                        message=(
                            f"Shape on {slide_note} may overflow: "
                            f"{total_chars} chars in {w_pt:.0f}×{h_pt:.0f}pt box."
                        ),
                    ))

        return issues

    # ── Brand compliance ───────────────────────────────────────────────────────

    def _validate_brand_compliance(
        self, deck: Deck, pptx_bytes: bytes
    ) -> list[ValidationIssue]:
        """
        Check rendered PPTX against brand rules (font names, colors).
        """
        issues: list[ValidationIssue] = []
        if not self._brand_rules:
            return issues

        try:
            prs = Presentation(io.BytesIO(pptx_bytes))
        except Exception:
            return issues

        # Build sets of allowed fonts
        allowed_heading_fonts: set[str] = set()
        allowed_body_fonts: set[str] = set()
        for rule in self._brand_rules:
            if rule.rule_type == "font" and "heading" in rule.name.lower():
                allowed_heading_fonts.update(v.lower() for v in rule.allowed_values)
            elif rule.rule_type == "font" and "body" in rule.name.lower():
                allowed_body_fonts.update(v.lower() for v in rule.allowed_values)

        # Scan all text runs for font violations
        theme_heading = deck.theme.fonts.heading.lower()
        theme_body = deck.theme.fonts.body.lower()
        all_allowed = (
            allowed_heading_fonts
            | allowed_body_fonts
            | {theme_heading, theme_body, "calibri", "calibri light"}
        )

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        fn = (run.font.name or "").lower()
                        if fn and fn not in all_allowed:
                            issues.append(ValidationIssue(
                                severity=ValidationSeverity.info,
                                code="UNEXPECTED_FONT",
                                message=(
                                    f"slide[{slide_idx + 1}]: font '{run.font.name}' "
                                    "is not in the approved font list."
                                ),
                            ))
                            break  # one per shape is enough

        return issues
