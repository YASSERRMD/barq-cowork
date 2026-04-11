"""
Main rendering pipeline: Deck JSON → PPTX bytes.

PresentationRenderer iterates over Deck.plan.slides, dispatches to the
appropriate SlideRenderer from the registry, embeds stable slide IDs in the
XML, and returns a RenderResult.
"""

from __future__ import annotations

import io
import os
from copy import deepcopy
from typing import Optional

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from lxml import etree

from .schema import (
    Deck,
    DeckTheme,
    RenderResult,
    Slide,
    SlideContent,
    SlideType,
)
from .slide_registry import SlideTypeRegistry

# ── Namespace for custom extension data ──────────────────────────────────────
_APP_NS = "http://barq.io/pptx/ext/2024"
_APP_TAG_ID = f"{{{_APP_NS}}}slideId"

# ── Slide layout index (blank layout) ─────────────────────────────────────────
_BLANK_LAYOUT_IDX = 6


class PresentationRenderer:
    """Render a Deck to a PPTX presentation."""

    def __init__(self, registry: SlideTypeRegistry):
        self._registry = registry

    # ── Public API ─────────────────────────────────────────────────────────────

    def render_deck(self, deck: Deck, output_path: str) -> RenderResult:
        """
        Render the complete deck to *output_path*.
        Returns a RenderResult with pptx_bytes and metadata.
        """
        prs = self._build_presentation(deck)
        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        with open(output_path, "wb") as f:
            f.write(pptx_bytes)

        return RenderResult(
            deck_id=deck.meta.id,
            version=deck.meta.version,
            pptx_bytes=pptx_bytes,
            slide_count=len(deck.plan.slides),
            output_path=output_path,
        )

    def render_to_bytes(self, deck: Deck) -> bytes:
        """Render deck entirely in memory and return raw PPTX bytes."""
        prs = self._build_presentation(deck)
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ── Internal build pipeline ────────────────────────────────────────────────

    def _build_presentation(self, deck: Deck) -> Presentation:
        prs = Presentation()
        # Set slide dimensions to widescreen 10 × 7.5 in
        prs.slide_width = Emu(9144000)
        prs.slide_height = Emu(6858000)

        for slide_model in deck.plan.slides:
            self._render_slide(prs, slide_model, deck.theme)

        return prs

    def _render_slide(self, prs: Presentation, slide_model: Slide, theme: DeckTheme) -> None:
        """Add one slide to prs and dispatch to the type renderer."""
        # Use blank layout (index 6) as base for all slides
        layout_idx = min(_BLANK_LAYOUT_IDX, len(prs.slide_layouts) - 1)
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Clear all default placeholder shapes added by the layout
        self._clear_slide_shapes(slide)

        # Inject heading override so renderers can access it
        content = slide_model.content.model_copy()
        object.__setattr__(content, "_heading_override", slide_model.heading)

        # Dispatch to renderer
        renderer = self._registry.get(slide_model.type)
        renderer.render(slide, content, theme)

        # Add speaker notes
        if slide_model.speaker_notes:
            self._set_speaker_notes(slide, slide_model.speaker_notes)

        # Embed stable ID into the slide XML
        self._embed_slide_id(slide, slide_model.id)

    def _clear_slide_shapes(self, slide) -> None:
        """Remove all shapes that were inherited from the layout."""
        sp_tree = slide.shapes._spTree
        for sp in list(sp_tree):
            tag = sp.tag.split("}")[-1] if "}" in sp.tag else sp.tag
            if tag in ("sp", "pic", "graphicFrame", "grpSp", "cxnSp"):
                sp_tree.remove(sp)

    def _embed_slide_id(self, slide, stable_id: str) -> None:
        """
        Store the stable slide ID in the slide's XML extLst.
        Uses a custom namespace element so it survives round-trips.
        The tag lives at:
            p:sld / p:cSld / p:extLst / p:ext / barq:slideId
        """
        cSld = slide._element.find(qn("p:cSld"))
        if cSld is None:
            return

        # Get or create extLst under cSld
        extLst = cSld.find(qn("p:extLst"))
        if extLst is None:
            extLst = etree.SubElement(cSld, qn("p:extLst"))

        # Remove any existing barq ext
        for ext in list(extLst):
            child = ext.find(_APP_TAG_ID)
            if child is not None:
                extLst.remove(ext)

        # Add new ext
        ext_el = etree.SubElement(extLst, qn("p:ext"))
        ext_el.set("uri", "{barq.io/slideId}")
        id_el = etree.SubElement(ext_el, _APP_TAG_ID)
        id_el.set("id", stable_id)

    def _read_slide_id(self, slide) -> Optional[str]:
        """Read the stable ID embedded in a slide's extLst. Returns None if absent."""
        cSld = slide._element.find(qn("p:cSld"))
        if cSld is None:
            return None
        extLst = cSld.find(qn("p:extLst"))
        if extLst is None:
            return None
        for ext in extLst:
            id_el = ext.find(_APP_TAG_ID)
            if id_el is not None:
                return id_el.get("id")
        return None

    def _set_speaker_notes(self, slide, notes_text: str) -> None:
        """Set the speaker notes for a slide."""
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = notes_text
        except Exception:
            pass

    def _apply_template_background(self, prs: Presentation, template_path: str) -> None:
        """
        Copy the slide master and layouts from an existing template PPTX.
        Used when a template_id maps to an actual file path.
        """
        if not os.path.exists(template_path):
            return
        try:
            template_prs = Presentation(template_path)
            # python-pptx does not expose a direct "copy master" API;
            # we simply load the template as the base presentation instead.
            # This method is a hook for subclasses that manage master XML directly.
        except Exception:
            pass

    def build_slide_id_map(self, pptx_bytes: bytes) -> dict[str, int]:
        """
        Parse PPTX bytes and return {stable_id: slide_index} mapping.
        Used by EditEngine to locate target slides by stable ID.
        """
        prs = Presentation(io.BytesIO(pptx_bytes))
        result: dict[str, int] = {}
        for idx, slide in enumerate(prs.slides):
            sid = self._read_slide_id(slide)
            if sid:
                result[sid] = idx
        return result
