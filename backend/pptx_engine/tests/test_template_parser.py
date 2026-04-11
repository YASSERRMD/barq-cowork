"""
Tests for TemplateParser.

Creates a minimal PPTX programmatically (no file dependency) and validates
that the parser extracts theme colors, fonts, layouts, and brand rules correctly.
"""

from __future__ import annotations

import io
import os
import sys
import tempfile

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

# Allow running tests from repo root without package install
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", ".."))

from pptx_engine.template_parser import TemplateParser
from pptx_engine.schema import ParsedTemplate, BrandRule


def _make_minimal_pptx() -> str:
    """Create a minimal PPTX file and return its path."""
    prs = Presentation()
    # Add a blank slide so the file is not empty
    prs.slides.add_slide(prs.slide_layouts[6])
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


@pytest.fixture(scope="module")
def minimal_pptx_path():
    path = _make_minimal_pptx()
    yield path
    try:
        os.unlink(path)
    except Exception:
        pass


@pytest.fixture(scope="module")
def parser():
    return TemplateParser()


# ── Tests ─────────────────────────────────────────────────────────────────────

class TestParsing:
    def test_parse_returns_parsed_template(self, parser, minimal_pptx_path):
        """parse() returns a ParsedTemplate instance."""
        result = parser.parse(minimal_pptx_path)
        assert isinstance(result, ParsedTemplate)

    def test_parse_slide_count(self, parser, minimal_pptx_path):
        """Parsed template reports correct slide count."""
        result = parser.parse(minimal_pptx_path)
        # We added 1 slide
        assert result.slide_count == 1

    def test_parse_slide_dimensions(self, parser, minimal_pptx_path):
        """Parsed template has non-zero slide dimensions."""
        result = parser.parse(minimal_pptx_path)
        assert result.slide_width > 0
        assert result.slide_height > 0

    def test_parse_source_path(self, parser, minimal_pptx_path):
        """ParsedTemplate records the source path."""
        result = parser.parse(minimal_pptx_path)
        assert result.source_path == minimal_pptx_path


class TestThemeColors:
    def test_theme_colors_extracted(self, parser, minimal_pptx_path):
        """Theme colors are extracted and produce a ThemeColors with valid hex values."""
        result = parser.parse(minimal_pptx_path)
        colors = result.theme_colors
        # All color fields should be non-empty hex strings
        for field in ["background", "surface", "accent", "text"]:
            val = getattr(colors, field)
            assert isinstance(val, str), f"{field} is not a string"
            assert val.startswith("#"), f"{field} does not start with #"
            assert len(val) in (4, 7), f"{field} has unexpected length: {val}"

    def test_theme_colors_are_valid_hex(self, parser, minimal_pptx_path):
        """All extracted color values are valid hex strings."""
        result = parser.parse(minimal_pptx_path)
        colors = result.theme_colors
        for field in type(colors).model_fields:
            val = getattr(colors, field)
            hex_part = val.lstrip("#")
            assert len(hex_part) in (3, 6), f"{field}={val} is not valid hex"
            # Should be convertible to int
            int(hex_part, 16)


class TestThemeFonts:
    def test_theme_fonts_extracted(self, parser, minimal_pptx_path):
        """Theme fonts are extracted with non-empty values."""
        result = parser.parse(minimal_pptx_path)
        fonts = result.theme_fonts
        assert fonts.heading, "Heading font should not be empty"
        assert fonts.body, "Body font should not be empty"

    def test_font_sizes_positive(self, parser, minimal_pptx_path):
        """Font size fields are positive integers."""
        result = parser.parse(minimal_pptx_path)
        fonts = result.theme_fonts
        assert fonts.size_heading > 0
        assert fonts.size_body > 0
        assert fonts.size_caption > 0


class TestLayouts:
    def test_layouts_extracted(self, parser, minimal_pptx_path):
        """At least one layout is extracted from a default PPTX."""
        result = parser.parse(minimal_pptx_path)
        assert len(result.layouts) > 0, "Should extract at least one layout"

    def test_layout_has_name(self, parser, minimal_pptx_path):
        """Each layout has a non-empty name."""
        result = parser.parse(minimal_pptx_path)
        for layout in result.layouts:
            assert layout.name, f"Layout missing name: {layout}"

    def test_layout_names_unique(self, parser, minimal_pptx_path):
        """Layout names are deduplicated."""
        result = parser.parse(minimal_pptx_path)
        names = [l.name for l in result.layouts]
        assert len(names) == len(set(names)), "Layout names should be unique"


class TestBrandRules:
    def test_generate_brand_rules_non_empty(self, parser, minimal_pptx_path):
        """generate_brand_rules returns at least one rule."""
        result = parser.parse(minimal_pptx_path)
        rules = parser.generate_brand_rules(result)
        assert len(rules) > 0, "Should generate at least one brand rule"

    def test_brand_rules_are_brand_rule_instances(self, parser, minimal_pptx_path):
        """All returned rules are BrandRule instances."""
        result = parser.parse(minimal_pptx_path)
        rules = parser.generate_brand_rules(result)
        for rule in rules:
            assert isinstance(rule, BrandRule)

    def test_brand_rules_have_names(self, parser, minimal_pptx_path):
        """All brand rules have non-empty names."""
        result = parser.parse(minimal_pptx_path)
        rules = parser.generate_brand_rules(result)
        for rule in rules:
            assert rule.name, f"Brand rule missing name: {rule}"

    def test_brand_rules_include_font_rule(self, parser, minimal_pptx_path):
        """There is at least one font-type brand rule."""
        result = parser.parse(minimal_pptx_path)
        rules = parser.generate_brand_rules(result)
        font_rules = [r for r in rules if r.rule_type == "font"]
        assert len(font_rules) > 0, "Should have at least one font brand rule"

    def test_brand_rules_include_content_rule(self, parser, minimal_pptx_path):
        """There is at least one content-type brand rule."""
        result = parser.parse(minimal_pptx_path)
        rules = parser.generate_brand_rules(result)
        content_rules = [r for r in rules if r.rule_type == "content"]
        assert len(content_rules) > 0, "Should have at least one content brand rule"
