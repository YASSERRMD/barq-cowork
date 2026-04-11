"""
Tests for PresentationRenderer.

Verifies that:
- render_deck produces valid PPTX bytes
- all slide types render without error
- stable IDs are embedded in output
- render_to_bytes is idempotent for same input
"""

from __future__ import annotations

import io
import os
import sys
import tempfile

import pytest
from pptx import Presentation

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", ".."))

from pptx_engine.renderer import PresentationRenderer
from pptx_engine.slide_registry import SlideTypeRegistry
from pptx_engine.schema import (
    CardItem,
    ChartSeries,
    ChartType,
    ChangeEntry,
    CompareColumn,
    Deck,
    DeckMeta,
    DeckPlan,
    DeckTheme,
    Slide,
    SlideContent,
    SlideID,
    SlideType,
    StatItem,
    TableData,
    TimelineItem,
)


# ── Fixtures ──────────────────────────────────────────────────────────────────

@pytest.fixture(scope="module")
def registry():
    return SlideTypeRegistry.build_default()


@pytest.fixture(scope="module")
def renderer(registry):
    return PresentationRenderer(registry)


def _make_deck(slides: list[Slide]) -> Deck:
    meta = DeckMeta(
        title="Test Deck",
        author="test",
        change_log=[
            ChangeEntry(version=1, author="test", action="create",
                        slide_ids_affected=[], description="test")
        ],
    )
    plan = DeckPlan(slides=slides)
    return Deck(meta=meta, theme=DeckTheme(), plan=plan)


def _all_slide_types() -> list[Slide]:
    """Create one slide of each supported type with representative content."""
    return [
        Slide(
            id=SlideID.generate(),
            type=SlideType.title,
            heading="Test Title",
            content=SlideContent(subtitle="A test subtitle"),
        ),
        Slide(
            id=SlideID.generate(),
            type=SlideType.bullets,
            heading="Bullet Points",
            content=SlideContent(points=["Point A", "Point B", "Point C"]),
        ),
        Slide(
            id=SlideID.generate(),
            type=SlideType.stats,
            heading="Key Metrics",
            content=SlideContent(stats=[
                StatItem(value="98%", label="Uptime"),
                StatItem(value="2.3M", label="Users"),
                StatItem(value="45%", label="Growth"),
                StatItem(value="$1.2M", label="Revenue"),
            ]),
        ),
        Slide(
            id=SlideID.generate(),
            type=SlideType.steps,
            heading="Process Steps",
            content=SlideContent(steps=["Define", "Design", "Build", "Test", "Deploy"]),
        ),
        Slide(
            id=SlideID.generate(),
            type=SlideType.cards,
            heading="Feature Cards",
            content=SlideContent(cards=[
                CardItem(icon="⚡", title="Speed", desc="Fast"),
                CardItem(icon="🔒", title="Security", desc="Secure"),
                CardItem(icon="📊", title="Analytics", desc="Insightful"),
            ]),
        ),
        Slide(
            id=SlideID.generate(),
            type=SlideType.chart,
            heading="Revenue Chart",
            content=SlideContent(
                chart_type=ChartType.column,
                chart_categories=["Q1", "Q2", "Q3", "Q4"],
                chart_series=[ChartSeries(name="Rev", values=[1.0, 2.0, 3.0, 4.0])],
            ),
        ),
        Slide(
            id=SlideID.generate(),
            type=SlideType.timeline,
            heading="Roadmap",
            content=SlideContent(timeline=[
                TimelineItem(date="Q1", title="Launch"),
                TimelineItem(date="Q2", title="Scale"),
                TimelineItem(date="Q3", title="Expand"),
            ]),
        ),
        Slide(
            id=SlideID.generate(),
            type=SlideType.compare,
            heading="Before vs After",
            content=SlideContent(
                left_column=CompareColumn(heading="Before", points=["Manual", "Slow"]),
                right_column=CompareColumn(heading="After", points=["Automated", "Fast"]),
            ),
        ),
        Slide(
            id=SlideID.generate(),
            type=SlideType.table,
            heading="Comparison Table",
            content=SlideContent(table=TableData(
                headers=["Plan", "Price", "Users"],
                rows=[["Basic", "$10", "5"], ["Pro", "$30", "25"]],
            )),
        ),
        Slide(
            id=SlideID.generate(),
            type=SlideType.blank,
            heading="Blank",
            content=SlideContent(body_text="Free-form content here."),
        ),
    ]


# ── Tests ─────────────────────────────────────────────────────────────────────

class TestRenderDeck:
    def test_render_deck_produces_bytes(self, renderer, tmp_path):
        """render_deck writes a non-empty PPTX file and returns RenderResult."""
        deck = _make_deck([
            Slide(id=SlideID.generate(), type=SlideType.title, heading="Hello",
                  content=SlideContent(subtitle="World"))
        ])
        out_path = str(tmp_path / "test.pptx")
        result = renderer.render_deck(deck, out_path)

        assert result.pptx_bytes is not None
        assert len(result.pptx_bytes) > 1000
        assert os.path.exists(out_path)

    def test_render_deck_produces_valid_pptx(self, renderer, tmp_path):
        """Output bytes are parseable as a valid PPTX."""
        deck = _make_deck([
            Slide(id=SlideID.generate(), type=SlideType.title, heading="Valid",
                  content=SlideContent(subtitle="Deck"))
        ])
        out_path = str(tmp_path / "valid.pptx")
        result = renderer.render_deck(deck, out_path)

        prs = Presentation(io.BytesIO(result.pptx_bytes))
        assert len(prs.slides) == 1

    def test_render_deck_slide_count_matches(self, renderer, tmp_path):
        """Rendered PPTX has the same number of slides as the deck plan."""
        slides = _all_slide_types()
        deck = _make_deck(slides)
        out_path = str(tmp_path / "all_types.pptx")
        result = renderer.render_deck(deck, out_path)

        prs = Presentation(io.BytesIO(result.pptx_bytes))
        assert len(prs.slides) == len(slides)


class TestAllSlideTypes:
    @pytest.mark.parametrize("slide", _all_slide_types(), ids=lambda s: s.type.value)
    def test_slide_type_renders_without_error(self, renderer, slide):
        """Each slide type renders without raising an exception."""
        deck = _make_deck([slide])
        pptx_bytes = renderer.render_to_bytes(deck)
        assert len(pptx_bytes) > 500

    @pytest.mark.parametrize("slide", _all_slide_types(), ids=lambda s: s.type.value)
    def test_slide_type_produces_valid_pptx(self, renderer, slide):
        """Each slide type produces parseable PPTX output."""
        deck = _make_deck([slide])
        pptx_bytes = renderer.render_to_bytes(deck)
        prs = Presentation(io.BytesIO(pptx_bytes))
        assert len(prs.slides) == 1


class TestStableIds:
    def test_stable_ids_embedded_in_output(self, renderer):
        """Stable IDs embedded during render can be read back from the PPTX."""
        slides = [
            Slide(id=SlideID.generate(), type=SlideType.title, heading="Slide 1",
                  content=SlideContent()),
            Slide(id=SlideID.generate(), type=SlideType.bullets, heading="Slide 2",
                  content=SlideContent(points=["A", "B"])),
        ]
        deck = _make_deck(slides)
        pptx_bytes = renderer.render_to_bytes(deck)

        # Read back IDs
        id_map = renderer.build_slide_id_map(pptx_bytes)
        assert len(id_map) == 2

        for slide in slides:
            assert slide.id in id_map, f"Slide ID {slide.id} not found in rendered PPTX"

    def test_stable_id_format_preserved(self, renderer):
        """Stable ID format 's-{8hex}' is preserved through render/read cycle."""
        slide_id = SlideID.generate()
        assert SlideID.is_valid(slide_id)

        slide = Slide(id=slide_id, type=SlideType.blank, heading="",
                      content=SlideContent())
        deck = _make_deck([slide])
        pptx_bytes = renderer.render_to_bytes(deck)

        id_map = renderer.build_slide_id_map(pptx_bytes)
        assert slide_id in id_map


class TestIdempotency:
    def test_render_to_bytes_same_deck_produces_same_slide_count(self, renderer):
        """render_to_bytes on the same deck produces same slide count."""
        deck = _make_deck([
            Slide(id=SlideID.generate(), type=SlideType.title, heading="Idempotent",
                  content=SlideContent(subtitle="Test"))
        ])
        bytes1 = renderer.render_to_bytes(deck)
        bytes2 = renderer.render_to_bytes(deck)

        prs1 = Presentation(io.BytesIO(bytes1))
        prs2 = Presentation(io.BytesIO(bytes2))
        assert len(prs1.slides) == len(prs2.slides)

    def test_render_produces_consistent_slide_structure(self, renderer):
        """Two renders of the same deck have matching IDs."""
        deck = _make_deck([
            Slide(id="s-aabbccdd", type=SlideType.bullets, heading="Consistent",
                  content=SlideContent(points=["X", "Y", "Z"]))
        ])
        bytes1 = renderer.render_to_bytes(deck)
        bytes2 = renderer.render_to_bytes(deck)

        map1 = renderer.build_slide_id_map(bytes1)
        map2 = renderer.build_slide_id_map(bytes2)

        assert set(map1.keys()) == set(map2.keys())
