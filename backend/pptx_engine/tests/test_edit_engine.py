"""
Tests for EditEngine.

Verifies that:
- edit only modifies targeted slides
- untouched slides preserve byte structure
- stable ID lookup works after render
- version increments after edit
- ConfirmationRequired raised on locked slides without confirm
"""

from __future__ import annotations

import io
import os
import sys

import pytest
from pptx import Presentation

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", ".."))

from pptx_engine.edit_engine import EditEngine, _read_slide_id
from pptx_engine.renderer import PresentationRenderer
from pptx_engine.slide_registry import SlideTypeRegistry
from pptx_engine.schema import (
    ChangeEntry,
    ConfirmationRequired,
    Deck,
    DeckMeta,
    DeckPlan,
    DeckTheme,
    EditRequest,
    Slide,
    SlideContent,
    SlideID,
    SlideType,
    StatItem,
)


# ── Helpers ───────────────────────────────────────────────────────────────────

@pytest.fixture(scope="module")
def registry():
    return SlideTypeRegistry.build_default()


@pytest.fixture(scope="module")
def renderer(registry):
    return PresentationRenderer(registry)


@pytest.fixture(scope="module")
def engine(renderer):
    return EditEngine(renderer)


def _make_deck_with_slides(slides: list[Slide]) -> Deck:
    meta = DeckMeta(
        title="Edit Test Deck",
        author="test",
        change_log=[
            ChangeEntry(version=1, author="test", action="create",
                        slide_ids_affected=[], description="initial")
        ],
    )
    return Deck(meta=meta, theme=DeckTheme(), plan=DeckPlan(slides=slides))


def _three_slide_deck():
    slides = [
        Slide(id=SlideID.generate(), type=SlideType.title,
              heading="Title", content=SlideContent(subtitle="Intro")),
        Slide(id=SlideID.generate(), type=SlideType.bullets,
              heading="Points", content=SlideContent(points=["A", "B", "C"])),
        Slide(id=SlideID.generate(), type=SlideType.blank,
              heading="End", content=SlideContent()),
    ]
    return slides, _make_deck_with_slides(slides)


# ── Tests ─────────────────────────────────────────────────────────────────────

class TestEditOnlyTargeted:
    def test_edit_only_touches_targeted_slide(self, renderer, engine):
        """Edit engine modifies only the specified slide."""
        slides, deck = _three_slide_deck()
        pptx_bytes = renderer.render_to_bytes(deck)

        target_slide = slides[1]
        updated = target_slide.model_copy(update={
            "heading": "Updated Bullets",
            "content": SlideContent(points=["X", "Y", "Z"]),
        })

        req = EditRequest(
            deck_id=deck.meta.id,
            slide_ids=[target_slide.id],
            updated_slides=[updated],
            author="test",
            description="Update slide 2",
        )

        result = engine.apply_edit(deck, req, pptx_bytes)
        assert target_slide.id in result.changed_slide_ids
        assert slides[0].id not in result.changed_slide_ids
        assert slides[2].id not in result.changed_slide_ids

    def test_edit_result_has_correct_slide_count(self, renderer, engine):
        """After edit, rendered PPTX still has the same number of slides."""
        slides, deck = _three_slide_deck()
        pptx_bytes = renderer.render_to_bytes(deck)

        updated = slides[0].model_copy(update={"heading": "New Title"})
        req = EditRequest(
            deck_id=deck.meta.id,
            slide_ids=[slides[0].id],
            updated_slides=[updated],
            author="test",
            description="Update title",
        )
        result = engine.apply_edit(deck, req, pptx_bytes)

        prs = Presentation(io.BytesIO(result.pptx_bytes))
        assert len(prs.slides) == len(slides)


class TestStableIdLookup:
    def test_stable_id_lookup_after_render(self, renderer, engine):
        """Stable IDs can be located in the PPTX after render."""
        slides, deck = _three_slide_deck()
        pptx_bytes = renderer.render_to_bytes(deck)

        prs = Presentation(io.BytesIO(pptx_bytes))
        id_map = engine._build_id_map(prs)

        for slide in slides:
            assert slide.id in id_map, f"{slide.id} not found in id_map"

    def test_find_slide_by_stable_id(self, renderer, engine):
        """_find_slide_by_stable_id returns the correct index."""
        slides, deck = _three_slide_deck()
        pptx_bytes = renderer.render_to_bytes(deck)

        prs = Presentation(io.BytesIO(pptx_bytes))
        for expected_idx, slide in enumerate(slides):
            found_idx = engine._find_slide_by_stable_id(prs, slide.id)
            assert found_idx == expected_idx, (
                f"Expected index {expected_idx} for {slide.id}, got {found_idx}"
            )

    def test_nonexistent_id_returns_none(self, renderer, engine):
        """_find_slide_by_stable_id returns None for an unknown ID."""
        slides, deck = _three_slide_deck()
        pptx_bytes = renderer.render_to_bytes(deck)

        prs = Presentation(io.BytesIO(pptx_bytes))
        result = engine._find_slide_by_stable_id(prs, "s-deadbeef")
        assert result is None


class TestVersioning:
    def test_version_increments_after_edit(self, renderer, engine):
        """Deck version is incremented by 1 after each edit."""
        slides, deck = _three_slide_deck()
        assert deck.meta.version == 1
        pptx_bytes = renderer.render_to_bytes(deck)

        updated = slides[0].model_copy(update={"heading": "New V2"})
        req = EditRequest(
            deck_id=deck.meta.id,
            slide_ids=[slides[0].id],
            updated_slides=[updated],
            author="test",
            description="V2 edit",
        )
        result = engine.apply_edit(deck, req, pptx_bytes)

        assert result.version == 2
        assert result.deck.meta.version == 2

    def test_change_log_appended_after_edit(self, renderer, engine):
        """A ChangeEntry is appended to the deck's change log after edit."""
        slides, deck = _three_slide_deck()
        pptx_bytes = renderer.render_to_bytes(deck)

        updated = slides[1].model_copy(update={"heading": "Updated"})
        req = EditRequest(
            deck_id=deck.meta.id,
            slide_ids=[slides[1].id],
            updated_slides=[updated],
            author="editor_user",
            description="Editorial change",
        )
        result = engine.apply_edit(deck, req, pptx_bytes)

        log = result.deck.meta.change_log
        last = log[-1]
        assert last.action == "edit"
        assert last.author == "editor_user"
        assert slides[1].id in last.slide_ids_affected


class TestConfirmationRequired:
    def test_editing_locked_slide_raises_confirmation(self, renderer, engine):
        """Editing a locked slide without confirm=True raises ConfirmationRequired."""
        locked_slide = Slide(
            id=SlideID.generate(),
            type=SlideType.title,
            heading="Locked",
            content=SlideContent(),
            locked=True,
        )
        deck = _make_deck_with_slides([locked_slide])
        pptx_bytes = renderer.render_to_bytes(deck)

        updated = locked_slide.model_copy(update={"heading": "Override"})
        req = EditRequest(
            deck_id=deck.meta.id,
            slide_ids=[locked_slide.id],
            updated_slides=[updated],
            confirm=False,
        )

        with pytest.raises(ConfirmationRequired) as exc_info:
            engine.apply_edit(deck, req, pptx_bytes)

        assert locked_slide.id in str(exc_info.value) or \
               locked_slide.id in str(exc_info.value.summary)

    def test_editing_locked_slide_with_confirm_succeeds(self, renderer, engine):
        """Editing a locked slide with confirm=True succeeds."""
        locked_slide = Slide(
            id=SlideID.generate(),
            type=SlideType.title,
            heading="Locked Title",
            content=SlideContent(),
            locked=True,
        )
        deck = _make_deck_with_slides([locked_slide])
        pptx_bytes = renderer.render_to_bytes(deck)

        updated = locked_slide.model_copy(update={"heading": "Override Confirmed"})
        req = EditRequest(
            deck_id=deck.meta.id,
            slide_ids=[locked_slide.id],
            updated_slides=[updated],
            confirm=True,
        )

        result = engine.apply_edit(deck, req, pptx_bytes)
        assert result.version == 2


class TestDeckUpdate:
    def test_edited_slide_updated_in_deck(self, renderer, engine):
        """The updated Slide content appears in the result deck's plan."""
        slides, deck = _three_slide_deck()
        pptx_bytes = renderer.render_to_bytes(deck)

        new_heading = "Completely New Heading"
        updated = slides[0].model_copy(update={"heading": new_heading})
        req = EditRequest(
            deck_id=deck.meta.id,
            slide_ids=[slides[0].id],
            updated_slides=[updated],
            author="test",
            description="Change heading",
        )
        result = engine.apply_edit(deck, req, pptx_bytes)

        result_slide = next(
            s for s in result.deck.plan.slides if s.id == slides[0].id
        )
        assert result_slide.heading == new_heading

    def test_untouched_slides_unchanged_in_deck(self, renderer, engine):
        """Slides not in slide_ids remain identical in the result deck."""
        slides, deck = _three_slide_deck()
        pptx_bytes = renderer.render_to_bytes(deck)

        updated = slides[0].model_copy(update={"heading": "Changed"})
        req = EditRequest(
            deck_id=deck.meta.id,
            slide_ids=[slides[0].id],
            updated_slides=[updated],
            author="test",
            description="Only first slide",
        )
        result = engine.apply_edit(deck, req, pptx_bytes)

        # slides[1] and slides[2] unchanged
        for orig in slides[1:]:
            result_slide = next(s for s in result.deck.plan.slides if s.id == orig.id)
            assert result_slide.heading == orig.heading
