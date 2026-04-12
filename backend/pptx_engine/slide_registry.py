"""
Slide type registry — one deterministic renderer per SlideType.

Each renderer receives a python-pptx Slide object and fills it with shapes
matching the dark glassmorphism aesthetic from gen_pptx.py:
  BG=#0F172A, CARD=#1E293B, ACCENT=theme.colors.accent
"""

from __future__ import annotations

import re
from typing import Protocol, runtime_checkable

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from lxml import etree

from .schema import (
    CardItem,
    ChartType,
    DeckTheme,
    SlideContent,
    SlideType,
    StatItem,
    TimelineItem,
)

# ── Canvas constants ──────────────────────────────────────────────────────────
W, H = 9144000, 6858000   # 10 × 7.5 inches in EMU

# ── Default dark palette ──────────────────────────────────────────────────────
_BG    = RGBColor(0x0F, 0x17, 0x2A)
_CARD  = RGBColor(0x1E, 0x29, 0x3B)
_CARD2 = RGBColor(0x0D, 0x1B, 0x2E)
_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
_OFF   = RGBColor(0xE2, 0xE8, 0xF0)
_MUTED = RGBColor(0x94, 0xA3, 0xB8)
_LINE  = RGBColor(0x2D, 0x3F, 0x55)


# ── Color helpers ─────────────────────────────────────────────────────────────

def _h(hex6: str) -> RGBColor:
    h = hex6.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _v(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def _lighter(c: RGBColor, amt: int = 60) -> RGBColor:
    return RGBColor(min(255, c[0] + amt), min(255, c[1] + amt), min(255, c[2] + amt))


def _darker(c: RGBColor, amt: int = 40) -> RGBColor:
    return RGBColor(max(0, c[0] - amt), max(0, c[1] - amt), max(0, c[2] - amt))


def _accent_from_theme(theme: DeckTheme) -> RGBColor:
    return _h(theme.colors.accent)


def _accent2_from_theme(theme: DeckTheme) -> RGBColor:
    return _h(theme.colors.accent2)


def _bg_from_theme(theme: DeckTheme) -> RGBColor:
    return _h(theme.colors.background)


def _card_from_theme(theme: DeckTheme) -> RGBColor:
    return _h(theme.colors.surface)


def _muted_from_theme(theme: DeckTheme) -> RGBColor:
    return _h(theme.colors.text_muted)


def _text_from_theme(theme: DeckTheme) -> RGBColor:
    return _h(theme.colors.text)


def _line_from_theme(theme: DeckTheme) -> RGBColor:
    return _h(theme.colors.border)



# ── Low-level XML primitives (mirror gen_pptx.py) ────────────────────────────

def _spPr(shape):
    return shape._element.find(qn("p:spPr"))


def _rm_fills(spPr):
    for t in ["a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"]:
        e = spPr.find(qn(t))
        if e is not None:
            spPr.remove(e)


def _rm_fills_el(el):
    for t in ["a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"]:
        e = el.find(qn(t))
        if e is not None:
            el.remove(e)


def xml_solid_alpha(shape, color: RGBColor, alpha_pct: float):
    spPr = _spPr(shape)
    _rm_fills(spPr)
    sf = etree.SubElement(spPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", _v(color))
    al = etree.SubElement(sc, qn("a:alpha"))
    al.set("val", str(int(alpha_pct * 1000)))


def xml_grad(shape, stops: list, angle_deg: float = 90):
    spPr = _spPr(shape)
    _rm_fills(spPr)
    gf = etree.SubElement(spPr, qn("a:gradFill"))
    gs_lst = etree.SubElement(gf, qn("a:gsLst"))
    for pos, color in stops:
        gs = etree.SubElement(gs_lst, qn("a:gs"))
        gs.set("pos", str(int(pos * 1000)))
        sc = etree.SubElement(gs, qn("a:srgbClr"))
        sc.set("val", _v(color))
    lin = etree.SubElement(gf, qn("a:lin"))
    lin.set("ang", str(int(angle_deg * 60000)))
    lin.set("scaled", "0")


def xml_shadow(shape, blur_pt: float = 8, dist_pt: float = 5,
               dir_deg: float = 315, alpha: int = 50):
    spPr = _spPr(shape)
    el = spPr.find(qn("a:effectLst"))
    if el is None:
        el = etree.SubElement(spPr, qn("a:effectLst"))
    old = el.find(qn("a:outerShdw"))
    if old is not None:
        el.remove(old)
    s = etree.SubElement(el, qn("a:outerShdw"))
    s.set("blurRad", str(int(blur_pt * 12700)))
    s.set("dist", str(int(dist_pt * 12700)))
    s.set("dir", str(int(dir_deg * 60000)))
    s.set("algn", "ctr")
    s.set("rotWithShape", "0")
    sc = etree.SubElement(s, qn("a:srgbClr"))
    sc.set("val", "000000")
    etree.SubElement(sc, qn("a:alpha")).set("val", str(int(alpha * 1000)))


def xml_stroke_alpha(shape, color: RGBColor, alpha_pct: float, w_pt: float = 1.0):
    spPr = _spPr(shape)
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(w_pt * 12700)))
    for c in list(ln):
        ln.remove(c)
    sf = etree.SubElement(ln, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", _v(color))
    etree.SubElement(sc, qn("a:alpha")).set("val", str(int(alpha_pct * 1000)))


def xml_round_corners(shape, adj: int = 25000):
    spPr = _spPr(shape)
    pg = spPr.find(qn("a:prstGeom"))
    if pg is None:
        return
    av = pg.find(qn("a:avLst"))
    if av is None:
        av = etree.SubElement(pg, qn("a:avLst"))
    for g in av.findall(qn("a:gd")):
        av.remove(g)
    gd = etree.SubElement(av, qn("a:gd"))
    gd.set("name", "adj")
    gd.set("fmla", f"val {adj}")


def xml_set_geom(shape, prst: str):
    spPr = _spPr(shape)
    pg = spPr.find(qn("a:prstGeom"))
    if pg is not None:
        pg.set("prst", prst)
        av = pg.find(qn("a:avLst"))
        if av is None:
            av = etree.SubElement(pg, qn("a:avLst"))
        for g in list(av):
            av.remove(g)


def _get_or_add(parent, tag: str):
    el = parent.find(tag)
    if el is None:
        el = etree.SubElement(parent, tag)
    return el


def _set_solid_fill(spPr, hex_val: str):
    _rm_fills_el(spPr)
    sf = etree.SubElement(spPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", hex_val)


def _set_no_border(spPr):
    ln = _get_or_add(spPr, qn("a:ln"))
    for c in list(ln):
        ln.remove(c)
    etree.SubElement(ln, qn("a:noFill"))


# ── Shape factories ───────────────────────────────────────────────────────────

def _solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _nofill(shape):
    shape.fill.background()


def _noline(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, color: RGBColor = None):
    s = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
    if color:
        _solid(s, color)
    else:
        _nofill(s)
    _noline(s)
    return s


def rrect(slide, x, y, w, h, color: RGBColor = None, adj: int = 22000):
    s = slide.shapes.add_shape(5, Emu(x), Emu(y), Emu(w), Emu(h))
    if color:
        _solid(s, color)
    else:
        _nofill(s)
    _noline(s)
    xml_round_corners(s, adj)
    return s


def ellipse(slide, x, y, w, h, color: RGBColor = None):
    s = slide.shapes.add_shape(9, Emu(x), Emu(y), Emu(w), Emu(h))
    if color:
        _solid(s, color)
    else:
        _nofill(s)
    _noline(s)
    return s


def txbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))


def set_para(shape, text: str, size_pt: float, color: RGBColor,
             bold: bool = False, font: str = "Calibri Light",
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
             italic: bool = False, wrap: bool = True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    p.text = ""
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def label(slide, x, y, w, h, text: str, size_pt: float, color: RGBColor,
          bold: bool = False, font: str = "Calibri Light",
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    s = txbox(slide, x, y, w, h)
    set_para(s, text, size_pt, color, bold, font, align, anchor)
    return s


# ── Chart dark-background helper (uses chart._element = chartSpace) ──────────

def _chart_dark_bg(chart, bg_color: RGBColor = None, area_color: RGBColor = None):
    if bg_color is None:
        bg_color = _CARD
    if area_color is None:
        area_color = _BG
    cs = chart._element  # chartSpace in python-pptx 1.x
    spPr = _get_or_add(cs, qn("c:spPr"))
    _set_solid_fill(spPr, _v(area_color))
    _set_no_border(spPr)

    chart_el = cs.find(qn("c:chart"))
    if chart_el is not None:
        plotArea = chart_el.find(qn("c:plotArea"))
        if plotArea is not None:
            spPr2 = _get_or_add(plotArea, qn("c:spPr"))
            _set_solid_fill(spPr2, _v(bg_color))
            _set_no_border(spPr2)


def _style_axis(axis, label_color: RGBColor = _MUTED, gridline_color: RGBColor = _LINE,
                hide_gridlines: bool = False):
    try:
        axis.tick_labels.font.color.rgb = label_color
        axis.tick_labels.font.size = Pt(9)
    except Exception:
        pass
    try:
        if hide_gridlines:
            axis.major_gridlines.format.line.fill.background()
        else:
            axis.major_gridlines.format.line.color.rgb = gridline_color
            axis.major_gridlines.format.line.width = Pt(0.5)
    except Exception:
        pass
    try:
        axis.format.line.fill.background()
    except Exception:
        pass


# ── Shared header / background ────────────────────────────────────────────────

def paint_bg(slide, accent: RGBColor, bg: RGBColor = None):
    if bg is None:
        bg = _BG
    rect(slide, 0, 0, W, H, bg)
    c1 = ellipse(slide, 6900000, -800000, 3000000, 3000000)
    xml_solid_alpha(c1, accent, 4)
    xml_stroke_alpha(c1, accent, 14, 2.5)
    c2 = ellipse(slide, 7500000, -300000, 1900000, 1900000)
    xml_solid_alpha(c2, accent, 3)
    xml_stroke_alpha(c2, accent, 8, 1.5)
    d = ellipse(slide, 180000, 6100000, 380000, 380000)
    xml_solid_alpha(d, accent, 10)


def paint_header(slide, heading: str, accent: RGBColor, font_name: str = "Calibri Light") -> int:
    """Draw gradient header bar + title. Returns content_top_y."""
    tb = rect(slide, 0, 0, W, 14000)
    xml_grad(tb, [(0, accent), (100, _darker(accent, 30))], angle_deg=0)
    t = txbox(slide, 457200, 55000, 7700000, 490000)
    set_para(t, heading, 31, _WHITE, bold=True, font=font_name,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    d = ellipse(slide, 380000, 277000, 55000, 55000, accent)
    rect(slide, 457200, 590000, 8229600, 8000, _LINE)
    return 660000  # content_top_y


def card_bg(slide, x, y, w, h, accent: RGBColor, adj: int = 22000, card_color: RGBColor = None):
    if card_color is None:
        card_color = _CARD
    c = rrect(slide, x, y, w, h, card_color, adj)
    xml_stroke_alpha(c, accent, 18, 1.0)
    xml_shadow(c, blur_pt=10, dist_pt=5, dir_deg=315, alpha=45)
    return c


# ── Semantic icon selection ───────────────────────────────────────────────────

# Maps keyword patterns to semantically relevant emoji icons
_ICON_MAP = [
    (["speed","fast","perform","quick","rapid","latency"],         "⚡"),
    (["secur","safe","protect","privacy","trust","compliance"],    "🔒"),
    (["integrat","connect","api","plugin","sync","webhook"],       "🔌"),
    (["data","analyt","insight","dashboard","metric","report"],    "📊"),
    (["growth","scale","expand","revenue","sales","profit"],       "📈"),
    (["ai","intelligence","smart","ml","neural","model","llm"],    "🧠"),
    (["cloud","infra","server","deploy","host","devops"],          "☁️"),
    (["team","people","talent","hire","human","collabor","user"],  "👥"),
    (["money","cost","financ","budget","invest","price","saving"], "💰"),
    (["global","world","internat","region","market","geo"],        "🌍"),
    (["innovat","pioneer","launch","new","creat","ideate"],        "🚀"),
    (["mobile","app","device","phone","tablet"],                   "📱"),
    (["time","schedul","deadline","calendar","sprint"],            "⏱️"),
    (["health","medical","patient","clinic","doctor","care"],      "🏥"),
    (["learn","educate","train","course","skill","certif"],        "📚"),
    (["automat","workflow","process","pipeline","bot","rpa"],      "⚙️"),
    (["support","help","customer","service","assist","success"],   "🎯"),
    (["green","sustain","eco","climate","carbon","renewable"],     "🌱"),
    (["award","quality","excel","best","top","leader","certif"],   "🏆"),
    (["search","discover","find","explor","research","audit"],     "🔍"),
    (["design","ui","ux","visual","brand","creative","art"],       "🎨"),
    (["document","report","write","content","publish"],            "📄"),
    (["network","connect","partner","ecosyst","alliance"],         "🤝"),
    (["robot","autonom","agent","bot","rpa"],                      "🤖"),
    (["shield","defend","threat","risk","mitigat","firewall"],     "🛡️"),
    (["notify","alert","monitor","observ","watch"],                "🔔"),
    (["code","develop","engineer","program","software","git"],     "💻"),
    (["test","qa","quality","debug","review","audit"],             "🧪"),
    (["key","access","auth","login","sso","identity"],             "🔑"),
    (["chart","graph","visual","trend","forecast"],                "📉"),
]


def _smart_icon(title: str, desc: str = "") -> str:
    text = (title + " " + desc).lower()
    for keywords, icon in _ICON_MAP:
        for kw in keywords:
            if kw in text:
                return icon
    # Generate a varied fallback based on title hash so same slide always gets same icon
    fallbacks = ["✦", "◆", "▸", "●", "★", "◉", "⬟", "⬡"]
    return fallbacks[hash(title) % len(fallbacks)]


# ── Protocol ──────────────────────────────────────────────────────────────────

@runtime_checkable
class SlideRenderer(Protocol):
    type: SlideType

    def render(self, slide, content: SlideContent, theme: DeckTheme) -> None:
        ...


# ══════════════════════════════════════════════════════════════════════════════
# RENDERERS
# ══════════════════════════════════════════════════════════════════════════════

class TitleSlideRenderer:
    type = SlideType.title

    def render(self, slide, content: SlideContent, theme: DeckTheme) -> None:
        accent = _accent_from_theme(theme)
        bg = _bg_from_theme(theme)
        font_h = theme.fonts.heading

        # Full dark background
        rect(slide, 0, 0, W, H, bg)

        # Large decorative gradient overlay
        overlay = rect(slide, 0, 0, W, H)
        xml_grad(overlay, [
            (0, _darker(accent, 60)),
            (50, bg),
            (100, bg),
        ], angle_deg=135)

        # Decorative circles
        c1 = ellipse(slide, W // 2, -400000, 4000000, 4000000)
        xml_solid_alpha(c1, accent, 5)
        xml_stroke_alpha(c1, accent, 20, 3)

        c2 = ellipse(slide, W // 2 + 1000000, 200000, 2200000, 2200000)
        xml_solid_alpha(c2, accent, 4)
        xml_stroke_alpha(c2, accent, 12, 1.5)

        # Bottom accent bar
        bar = rect(slide, 0, H - 30000, W, 30000)
        xml_grad(bar, [(0, accent), (50, _lighter(accent, 40)), (100, accent)], 0)

        # Heading
        heading_text = getattr(content, "_heading_override", "") or ""
        t = txbox(slide, 457200, 1800000, W - 914400, 1400000)
        set_para(t, heading_text, 52, _WHITE, bold=True, font=font_h,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Subtitle
        subtitle = content.subtitle or ""
        if subtitle:
            s = txbox(slide, 914400, 3400000, W - 1828800, 900000)
            set_para(s, subtitle, 20, _OFF, font=theme.fonts.body,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Decorative line below title
        ln = rect(slide, W // 2 - 1000000, 3300000, 2000000, 8000)
        xml_grad(ln, [(0, bg), (50, accent), (100, bg)], 0)


class BulletsSlideRenderer:
    type = SlideType.bullets

    def render(self, slide, content: SlideContent, theme: DeckTheme) -> None:
        accent = _accent_from_theme(theme)
        bg = _bg_from_theme(theme)
        card = _card_from_theme(theme)
        paint_bg(slide, accent, bg)
        top = paint_header(slide, getattr(content, "_heading_override", ""), accent,
                           theme.fonts.heading)

        points = content.points or []
        if not points:
            return

        n = len(points)
        avail = H - top - 180000
        two_col = n >= 5

        if two_col:
            cols, col_w, col_xs = 2, 3962400, [457200, 4724400]
            rows = (n + 1) // 2
        else:
            cols, col_w, col_xs = 1, 8229600, [457200]
            rows = n

        gap = 90000
        ch = min((avail - (rows - 1) * gap) // rows, 1450000)
        total = rows * ch + (rows - 1) * gap
        sy = top + (avail - total) // 2

        for i, pt in enumerate(points):
            col, row = i % cols, i // cols
            cx, cy = col_xs[col], sy + row * (ch + gap)
            cw = col_w

            card_bg(slide, cx, cy, cw, ch, accent, card_color=card)

            strip = rect(slide, cx, cy + ch // 8, 11000, ch * 3 // 4)
            xml_grad(strip, [(0, _lighter(accent, 60)), (50, accent),
                             (100, _darker(accent, 20))], 90)

            ns = min(int(ch * 0.52), 440000)
            nx, ny = cx + 190000, cy + (ch - ns) // 2
            nb = ellipse(slide, nx, ny, ns, ns)
            xml_grad(nb, [(0, _lighter(accent, 50)), (100, accent)], 135)
            xml_shadow(nb, 4, 2, 315, 45)
            label(slide, nx, ny, ns, ns, str(i + 1),
                  13 if ns > 350000 else 11, _WHITE, bold=True, font="Calibri")

            tx = nx + ns + 160000
            tw = cw - (tx - cx) - 120000
            t = txbox(slide, tx, cy + 55000, tw, ch - 110000)
            set_para(t, pt, 15 if ch > 900000 else 12, _OFF,
                     font=theme.fonts.body, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


class StatsSlideRenderer:
    type = SlideType.stats

    def render(self, slide, content: SlideContent, theme: DeckTheme) -> None:
        accent = _accent_from_theme(theme)
        bg = _bg_from_theme(theme)
        card = _card_from_theme(theme)
        paint_bg(slide, accent, bg)
        top = paint_header(slide, getattr(content, "_heading_override", ""), accent,
                           theme.fonts.heading)

        items = list(content.stats or [])
        if not items and content.points:
            for p in content.points:
                m = re.match(r"^([^\s]+)\s+(.*)", p.strip())
                items.append(StatItem(
                    value=m.group(1) if m else p,
                    label=m.group(2) if m else "",
                ))
        items = items[:4]
        if not items:
            return

        n = len(items)
        avail_w = 8229600
        gap = 200000
        cw = (avail_w - (n - 1) * gap) // n
        ch = 2600000
        sy = top + 60000
        sx = 457200
        numeric_vals: list[float] = []

        for i, it in enumerate(items):
            cx = sx + i * (cw + gap)
            card_bg(slide, cx, sy, cw, ch, accent, card_color=card)
            ts = rect(slide, cx, sy, cw, ch // 9)
            xml_grad(ts, [(0, accent), (100, _lighter(accent, 50))], 0)

            val = it.value
            label(slide, cx + 15000, sy + ch // 9 + ch // 18, cw - 30000, ch * 5 // 14,
                  val, 46, accent, bold=True, font=theme.fonts.heading)

            pct_m = re.search(r"(\d+(?:\.\d+)?)\s*%", val)
            bar_y = sy + ch // 9 + ch * 5 // 14 + ch // 14
            if pct_m:
                pct = float(pct_m.group(1)) / 100.0
                rrect(slide, cx + 50000, bar_y, cw - 100000, ch // 14, _LINE, adj=50000)
                bw = int((cw - 100000) * min(pct, 1.0))
                if bw > 5000:
                    fb = rrect(slide, cx + 50000, bar_y, bw, ch // 14, accent, adj=50000)
                    xml_grad(fb, [(0, accent), (100, _lighter(accent, 60))], 0)
                lbl_y = bar_y + ch // 14 + ch // 20
            else:
                lbl_y = bar_y

            label(slide, cx + 15000, lbl_y, cw - 30000, ch // 6,
                  it.label, 14, _WHITE, bold=True, font=theme.fonts.body)
            if it.desc:
                label(slide, cx + 15000, lbl_y + ch // 6, cw - 30000, ch // 8,
                      it.desc, 11, _MUTED, font=theme.fonts.body)

            nm = re.search(r"[\d,.]+", val)
            numeric_vals.append(float(nm.group().replace(",", "")) if nm else 0)

        # Bar chart below KPIs
        chart_top = sy + ch + 200000
        chart_h = H - chart_top - 200000
        if chart_h < 600000:
            return

        cd = CategoryChartData()
        cd.categories = [it.label or it.value for it in items]
        cd.add_series("Value", numeric_vals)

        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Emu(457200), Emu(chart_top), Emu(8229600), Emu(chart_h), cd,
        )
        ch_obj = gf.chart
        ch_obj.has_legend = False
        _chart_dark_bg(ch_obj, card, bg)

        try:
            ser = ch_obj.series[0]
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = accent
            ser.data_labels.show_value = True
            ser.data_labels.font.color.rgb = _WHITE
            ser.data_labels.font.size = Pt(10)
            ser.data_labels.font.bold = True
        except Exception:
            pass
        try:
            _style_axis(ch_obj.category_axis, _MUTED, _LINE)
            _style_axis(ch_obj.value_axis, _MUTED, _LINE)
        except Exception:
            pass


class StepsSlideRenderer:
    type = SlideType.steps

    def render(self, slide, content: SlideContent, theme: DeckTheme) -> None:
        accent = _accent_from_theme(theme)
        bg = _bg_from_theme(theme)
        paint_bg(slide, accent, bg)
        top = paint_header(slide, getattr(content, "_heading_override", ""), accent,
                           theme.fonts.heading)

        points = content.steps or content.points or []
        if not points:
            return

        n = min(len(points), 6)
        two_rows = n >= 5
        rows_data = [points[:3], points[3:n]] if two_rows else [points[:n]]

        avail_h = H - top - 200000
        row_h = avail_h // len(rows_data)
        step_idx = 0

        for row_i, row_pts in enumerate(rows_data):
            rn = len(row_pts)
            chevron_w = (8229600 - 60000 * (rn - 1)) // rn
            chevron_h = int(row_h * 0.48)
            cy = top + row_i * row_h + (row_h - chevron_h) // 2 - 40000
            text_y = cy + chevron_h + 80000
            text_h = row_h - chevron_h - 100000

            for i, pt in enumerate(row_pts):
                step_idx += 1
                cx = 457200 + i * (chevron_w + 60000)
                geom = "home" if i == 0 else "chevron"
                shp = rect(slide, cx, cy, chevron_w, chevron_h)
                xml_grad(shp, [(0, _lighter(accent, 30)), (100, _darker(accent, 20))], 135)
                xml_set_geom(shp, geom)
                xml_shadow(shp, 8, 4, 315, 40)

                ns = min(chevron_h // 3, 380000)
                nx = cx + (chevron_w - ns) // 2
                ny = cy - ns // 2
                nb = ellipse(slide, nx, ny, ns, ns)
                xml_grad(nb, [(0, _WHITE), (100, _OFF)], 135)
                xml_shadow(nb, 3, 2, 315, 30)
                label(slide, nx, ny, ns, ns, str(step_idx),
                      12, accent, bold=True, font="Calibri")

                inner_w = int(chevron_w * (0.7 if i < rn - 1 else 0.85))
                t = txbox(slide, cx + 20000, cy + 20000, inner_w, chevron_h - 40000)
                short = pt[:45] + "…" if len(pt) > 45 else pt
                set_para(t, short, 12, _WHITE, font=theme.fonts.body,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

                if len(pt) > 10:
                    t2 = txbox(slide, cx, text_y, chevron_w, text_h)
                    set_para(t2, pt, 11, _MUTED, font=theme.fonts.body,
                             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


class CardsSlideRenderer:
    type = SlideType.cards

    def render(self, slide, content: SlideContent, theme: DeckTheme) -> None:
        accent = _accent_from_theme(theme)
        accent2 = _accent2_from_theme(theme)
        bg = _bg_from_theme(theme)
        card_clr = _card_from_theme(theme)
        paint_bg(slide, accent, bg)
        top = paint_header(slide, getattr(content, "_heading_override", ""), accent,
                           theme.fonts.heading)

        cards = content.cards or []
        if not cards and content.points:
            cards = [CardItem(icon="●", title=p[:40], desc="") for p in content.points]

        cards = cards[:6]
        if not cards:
            return

        n = len(cards)
        cols = 3 if n >= 4 else min(n, 3)
        rows = (n + cols - 1) // cols
        avail_w = 8229600
        avail_h = H - top - 200000
        gap_x, gap_y = 200000, 200000
        cw = (avail_w - (cols - 1) * gap_x) // cols
        ch = (avail_h - (rows - 1) * gap_y) // rows
        sx, sy = 457200, top + 80000

        accent_cycle = [accent, accent2, _lighter(accent, 30)]

        for i, card in enumerate(cards):
            col = i % cols
            row = i // cols
            cx = sx + col * (cw + gap_x)
            cy = sy + row * (ch + gap_y)
            ac = accent_cycle[i % len(accent_cycle)]

            # Card background
            c = card_bg(slide, cx, cy, cw, ch, ac, card_color=card_clr)

            # Top accent bar
            tb = rect(slide, cx, cy, cw, ch // 10)
            xml_grad(tb, [(0, ac), (100, _lighter(ac, 40))], 0)

            # Icon circle
            icon_r = min(ch // 5, 500000)
            ix = cx + (cw - icon_r) // 2
            iy = cy + ch // 8
            ic = ellipse(slide, ix, iy, icon_r, icon_r)
            xml_grad(ic, [(0, _lighter(ac, 40)), (100, ac)], 135)
            xml_shadow(ic, 4, 2, 315, 40)
            # Use smart icon selection: replace missing or generic fallback icons
            _generic_icons = {"★", "●", "◈", "⬡", "◎", "✦", "◆", "▸", "◉", "⬟"}
            raw_icon = card.icon if card.icon else ""
            if not raw_icon or raw_icon in _generic_icons:
                icon_text = _smart_icon(card.title, card.desc)
            else:
                icon_text = raw_icon
            label(slide, ix, iy, icon_r, icon_r, icon_text,
                  min(18, icon_r // 30000), _WHITE, font="Segoe UI Symbol")

            # Title
            title_y = iy + icon_r + 80000
            title_h = ch // 5
            label(slide, cx + 40000, title_y, cw - 80000, title_h,
                  card.title, 14, _WHITE, bold=True, font=theme.fonts.body,
                  align=PP_ALIGN.CENTER)

            # Description
            if card.desc:
                desc_y = title_y + title_h + 30000
                desc_h = ch - (desc_y - cy) - 80000
                if desc_h > 100000:
                    t = txbox(slide, cx + 40000, desc_y, cw - 80000, desc_h)
                    set_para(t, card.desc, 11, _MUTED, font=theme.fonts.body,
                             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


class ChartSlideRenderer:
    type = SlideType.chart

    _CHART_TYPE_MAP = {
        ChartType.column: XL_CHART_TYPE.COLUMN_CLUSTERED,
        ChartType.bar: XL_CHART_TYPE.BAR_CLUSTERED,
        ChartType.line: XL_CHART_TYPE.LINE,
        ChartType.pie: XL_CHART_TYPE.PIE,
        ChartType.doughnut: XL_CHART_TYPE.DOUGHNUT,
        ChartType.area: XL_CHART_TYPE.AREA,
        ChartType.scatter: XL_CHART_TYPE.XY_SCATTER,
    }

    def render(self, slide, content: SlideContent, theme: DeckTheme) -> None:
        accent = _accent_from_theme(theme)
        bg = _bg_from_theme(theme)
        card = _card_from_theme(theme)
        paint_bg(slide, accent, bg)
        top = paint_header(slide, getattr(content, "_heading_override", ""), accent,
                           theme.fonts.heading)

        chart_type_enum = self._CHART_TYPE_MAP.get(
            content.chart_type or ChartType.column,
            XL_CHART_TYPE.COLUMN_CLUSTERED,
        )

        categories = content.chart_categories or ["Q1", "Q2", "Q3", "Q4"]
        series_list = content.chart_series or []

        cd = CategoryChartData()
        cd.categories = categories

        accent_colors = [accent, _accent2_from_theme(theme), _lighter(accent, 40)]

        if series_list:
            for s in series_list:
                cd.add_series(s.name, [float(v) for v in s.values])
        else:
            cd.add_series("Series 1", [25.0, 40.0, 35.0, 55.0][:len(categories)])

        chart_margin = 200000
        gf = slide.shapes.add_chart(
            chart_type_enum,
            Emu(457200), Emu(top + chart_margin),
            Emu(8229600), Emu(H - top - chart_margin * 2 - 100000),
            cd,
        )
        ch_obj = gf.chart
        _chart_dark_bg(ch_obj, card, bg)

        muted = _muted_from_theme(theme)
        text_color = _text_from_theme(theme)
        line_color = _line_from_theme(theme)

        if content.chart_title:
            ch_obj.has_title = True
            ch_obj.chart_title.text_frame.text = content.chart_title
            try:
                ch_obj.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = text_color
                ch_obj.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
                ch_obj.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
            except Exception:
                pass

        # Style series with theme accent colors
        for i, ser in enumerate(ch_obj.series):
            try:
                clr = accent_colors[i % len(accent_colors)]
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = clr
            except Exception:
                pass

        ch_obj.has_legend = len(series_list) > 1
        if ch_obj.has_legend:
            ch_obj.legend.position = XL_LEGEND_POSITION.BOTTOM
            try:
                ch_obj.legend.font.color.rgb = muted
                ch_obj.legend.font.size = Pt(10)
            except Exception:
                pass

        try:
            _style_axis(ch_obj.category_axis, muted, line_color)
            _style_axis(ch_obj.value_axis, muted, line_color)
        except Exception:
            pass

        # Y-axis label if provided
        if content.y_label:
            try:
                ch_obj.value_axis.axis_title.text_frame.text = content.y_label
                ch_obj.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.color.rgb = muted
                ch_obj.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
                ch_obj.value_axis.has_title = True
            except Exception:
                pass


class TimelineSlideRenderer:
    type = SlideType.timeline

    def render(self, slide, content: SlideContent, theme: DeckTheme) -> None:
        accent = _accent_from_theme(theme)
        bg = _bg_from_theme(theme)
        paint_bg(slide, accent, bg)
        top = paint_header(slide, getattr(content, "_heading_override", ""), accent,
                           theme.fonts.heading)

        items = content.timeline or []
        if not items and content.points:
            items = [TimelineItem(date=f"Phase {i+1}", title=p)
                     for i, p in enumerate(content.points)]

        items = items[:7]
        if not items:
            return

        n = len(items)
        avail_h = H - top - 300000
        spine_y = top + avail_h // 2
        sx = 457200
        ex = W - 457200
        spine_w = ex - sx
        node_gap = spine_w // n
        spine_h = 16000

        # Spine
        s = rect(slide, sx, spine_y, spine_w, spine_h)
        xml_grad(s, [(0, _darker(accent, 20)), (50, accent), (100, _darker(accent, 20))], 0)

        for i, item in enumerate(items):
            cx = sx + i * node_gap + node_gap // 2
            is_top = i % 2 == 0
            node_r = 200000
            nx = cx - node_r // 2
            ny = spine_y + spine_h // 2 - node_r // 2

            nb = ellipse(slide, nx, ny, node_r, node_r)
            xml_grad(nb, [(0, _lighter(accent, 40)), (100, accent)], 135)
            xml_shadow(nb, 4, 2, 315, 50)
            xml_stroke_alpha(nb, _WHITE, 30, 1.5)

            dot_num_size = min(node_r, 180000)
            label(slide, nx, ny, node_r, node_r, str(i + 1),
                  10, _WHITE, bold=True, font="Calibri")

            # Connector line
            line_h = 600000
            lx = cx - 8000
            if is_top:
                ly = ny - line_h
            else:
                ly = ny + node_r
            ln = rect(slide, lx, ly, 16000, line_h)
            xml_solid_alpha(ln, accent, 40)

            # Date label
            date_y = ly - 250000 if is_top else ly + line_h + 30000
            label(slide, cx - node_gap // 2, date_y, node_gap, 230000,
                  item.date, 10, accent, bold=True, font=theme.fonts.body,
                  align=PP_ALIGN.CENTER)

            # Title
            title_y = date_y + 230000 if is_top else date_y + 230000
            if not is_top:
                title_y = ly + line_h + 60000
            label(slide, cx - node_gap // 2, title_y, node_gap, 300000,
                  item.title, 12, _WHITE, bold=True, font=theme.fonts.body,
                  align=PP_ALIGN.CENTER)

            # Desc
            if item.desc:
                desc_y = title_y + 300000
                t = txbox(slide, cx - node_gap // 2, desc_y, node_gap, 300000)
                set_para(t, item.desc, 10, _MUTED, font=theme.fonts.body,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


class CompareSlideRenderer:
    type = SlideType.compare

    def render(self, slide, content: SlideContent, theme: DeckTheme) -> None:
        accent = _accent_from_theme(theme)
        accent2 = _accent2_from_theme(theme)
        bg = _bg_from_theme(theme)
        card = _card_from_theme(theme)
        paint_bg(slide, accent, bg)
        top = paint_header(slide, getattr(content, "_heading_override", ""), accent,
                           theme.fonts.heading)

        left = content.left_column
        right = content.right_column
        if left is None:
            left_points = (content.points or [])[:len(content.points or []) // 2]
            left = type("Col", (), {
                "heading": "Option A",
                "points": left_points,
                "color": None,
            })()
        if right is None:
            half = len(content.points or []) // 2
            right_points = (content.points or [])[half:]
            right = type("Col", (), {
                "heading": "Option B",
                "points": right_points,
                "color": None,
            })()

        col_w = (W - 914400 - 200000) // 2
        avail_h = H - top - 200000
        left_x = 457200
        right_x = 457200 + col_w + 200000
        col_top = top + 100000

        for col_x, col, col_accent in [
            (left_x, left, _h(left.color) if getattr(left, "color", None) else accent),
            (right_x, right, _h(right.color) if getattr(right, "color", None) else accent2),
        ]:
            # Column card
            c = card_bg(slide, col_x, col_top, col_w, avail_h, col_accent, card_color=card)
            # Header bar
            hb = rect(slide, col_x, col_top, col_w, 380000)
            xml_grad(hb, [(0, col_accent), (100, _darker(col_accent, 30))], 0)
            xml_round_corners(hb, 22000)

            label(slide, col_x, col_top, col_w, 380000,
                  getattr(col, "heading", ""), 18, _WHITE, bold=True,
                  font=theme.fonts.heading, align=PP_ALIGN.CENTER)

            # Points
            pts = getattr(col, "points", []) or []
            pt_top = col_top + 400000
            pt_h = min((avail_h - 420000) // max(len(pts), 1), 700000)

            for j, pt in enumerate(pts[:8]):
                py = pt_top + j * pt_h
                # Check mark circle
                ck = ellipse(slide, col_x + 60000, py + pt_h // 2 - 120000,
                             240000, 240000)
                xml_grad(ck, [(0, col_accent), (100, _darker(col_accent, 20))], 135)
                label(slide, col_x + 60000, py + pt_h // 2 - 120000,
                      240000, 240000, "✓", 12, _WHITE, bold=True, font="Segoe UI Symbol")

                t = txbox(slide, col_x + 340000, py + 40000,
                          col_w - 400000, pt_h - 80000)
                set_para(t, pt, 13, _OFF, font=theme.fonts.body,
                         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


class TableSlideRenderer:
    type = SlideType.table

    def render(self, slide, content: SlideContent, theme: DeckTheme) -> None:
        accent = _accent_from_theme(theme)
        bg = _bg_from_theme(theme)
        card = _card_from_theme(theme)
        card2 = _darker(card, 20)
        paint_bg(slide, accent, bg)
        top = paint_header(slide, getattr(content, "_heading_override", ""), accent,
                           theme.fonts.heading)

        table_data = content.table
        if table_data is None:
            return

        headers = table_data.headers or []
        rows = table_data.rows or []
        if not headers:
            return

        ncols = len(headers)
        nrows = len(rows)

        table_top = top + 100000
        table_h = H - table_top - 200000
        table_w = 8229600

        row_h = min(table_h // (nrows + 1), 800000)
        header_h = int(row_h * 1.3)
        col_w = table_w // ncols

        # Header row
        for j, hdr in enumerate(headers):
            hx = 457200 + j * col_w
            hc = rect(slide, hx, table_top, col_w, header_h)
            xml_grad(hc, [(0, accent), (100, _darker(accent, 30))], angle_deg=90)
            if j > 0:
                xml_stroke_alpha(hc, _LINE, 30, 1.0)
            label(slide, hx, table_top, col_w, header_h, hdr,
                  13, _WHITE, bold=True, font=theme.fonts.body, align=PP_ALIGN.CENTER)

        # Data rows
        for i, row in enumerate(rows[:12]):
            ry = table_top + header_h + i * row_h
            row_color = card if i % 2 == 0 else card2

            for j, cell in enumerate(row[:ncols]):
                cx = 457200 + j * col_w
                rc = rect(slide, cx, ry, col_w, row_h, row_color)
                xml_stroke_alpha(rc, _LINE, 20, 0.5)

                t = txbox(slide, cx + 60000, ry + 20000, col_w - 120000, row_h - 40000)
                set_para(t, str(cell), 12, _OFF, font=theme.fonts.body,
                         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


class BlankSlideRenderer:
    type = SlideType.blank

    def render(self, slide, content: SlideContent, theme: DeckTheme) -> None:
        accent = _accent_from_theme(theme)
        bg = _bg_from_theme(theme)
        paint_bg(slide, accent, bg)

        heading = getattr(content, "_heading_override", "")
        if heading:
            paint_header(slide, heading, accent, theme.fonts.heading)

        if content.body_text:
            t = txbox(slide, 457200, 800000, W - 914400, H - 1000000)
            set_para(t, content.body_text, 16, _OFF, font=theme.fonts.body,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# REGISTRY
# ══════════════════════════════════════════════════════════════════════════════

class SlideTypeRegistry:
    """Maps SlideType → SlideRenderer instances."""

    def __init__(self):
        self._renderers: dict[SlideType, SlideRenderer] = {}

    def register(self, renderer: SlideRenderer) -> None:
        self._renderers[renderer.type] = renderer

    def get(self, slide_type: SlideType) -> SlideRenderer:
        renderer = self._renderers.get(slide_type)
        if renderer is None:
            raise KeyError(f"No renderer registered for slide type: {slide_type}")
        return renderer

    def has(self, slide_type: SlideType) -> bool:
        return slide_type in self._renderers

    @staticmethod
    def build_default() -> "SlideTypeRegistry":
        """Create a registry with all built-in renderers registered."""
        reg = SlideTypeRegistry()
        reg.register(TitleSlideRenderer())
        reg.register(BulletsSlideRenderer())
        reg.register(StatsSlideRenderer())
        reg.register(StepsSlideRenderer())
        reg.register(CardsSlideRenderer())
        reg.register(ChartSlideRenderer())
        reg.register(TimelineSlideRenderer())
        reg.register(CompareSlideRenderer())
        reg.register(TableSlideRenderer())
        reg.register(BlankSlideRenderer())
        return reg
