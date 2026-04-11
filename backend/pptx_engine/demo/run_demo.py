#!/usr/bin/env python3
"""
End-to-end demo for the Barq PPTX Engine.

Demonstrates the complete pipeline without a running server:
1. Create an "AI Strategy 2026" deck with 8 slides
2. Run validation
3. Edit slides 2 and 5 with new content
4. Verify no-touch preservation
5. Export final PPTX to /tmp/demo_deck.pptx
6. Print the change log
"""

from __future__ import annotations

import io
import os
import sys
import tempfile
from datetime import datetime

# Allow running from any directory
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", ".."))

from pptx import Presentation

from pptx_engine.deck_store import DeckStore
from pptx_engine.edit_engine import EditEngine
from pptx_engine.planner import create_plan
from pptx_engine.renderer import PresentationRenderer
from pptx_engine.schema import (
    CardItem,
    ChangeEntry,
    ChartSeries,
    ChartType,
    CompareColumn,
    CreateDeckRequest,
    Deck,
    DeckMeta,
    DeckPlan,
    DeckTheme,
    EditRequest,
    PlanRequest,
    Slide,
    SlideContent,
    SlideID,
    SlideType,
    StatItem,
    TimelineItem,
    ThemeColors,
    ThemeFonts,
)
from pptx_engine.slide_registry import SlideTypeRegistry
from pptx_engine.validator import DeckValidator

# ── Helpers ───────────────────────────────────────────────────────────────────

def _sep(title: str = ""):
    width = 72
    if title:
        pad = (width - len(title) - 2) // 2
        print(f"\n{'─' * pad} {title} {'─' * pad}")
    else:
        print("─" * width)


def _print_deck_summary(deck: Deck):
    print(f"  Title    : {deck.meta.title}")
    print(f"  ID       : {deck.meta.id}")
    print(f"  Author   : {deck.meta.author}")
    print(f"  Version  : {deck.meta.version}")
    print(f"  Slides   : {len(deck.plan.slides)}")
    print(f"  Duration : ~{deck.plan.estimated_duration_min:.1f} min")
    for i, s in enumerate(deck.plan.slides):
        locked = " [LOCKED]" if s.locked else ""
        print(f"    [{i+1:02d}] {s.type.value:12s} | {s.id} | {s.heading[:45]}{locked}")


def _print_change_log(change_log: list[ChangeEntry]):
    if not change_log:
        print("  (empty change log)")
        return
    for entry in change_log:
        ts = entry.timestamp.strftime("%Y-%m-%d %H:%M:%S")
        affected = ", ".join(entry.slide_ids_affected) if entry.slide_ids_affected else "—"
        print(f"  v{entry.version:02d} | {ts} | {entry.author:12s} | {entry.action:14s} | {entry.description[:40]}")
        print(f"       affected: {affected}")


# ── Demo deck definition ──────────────────────────────────────────────────────

def _build_ai_strategy_deck() -> Deck:
    """Build an 'AI Strategy 2026' deck with 8 diverse slides."""
    slides = [
        # 1. Title
        Slide(
            id=SlideID.generate(),
            type=SlideType.title,
            heading="AI Strategy 2026",
            content=SlideContent(
                subtitle="Building the Intelligent Enterprise"
            ),
            speaker_notes="Welcome the audience. Introduce the AI transformation vision.",
        ),
        # 2. Overview / Bullets  ← will be EDITED
        Slide(
            id=SlideID.generate(),
            type=SlideType.bullets,
            heading="Strategic Pillars",
            content=SlideContent(
                points=[
                    "Automate repetitive processes with GenAI",
                    "Embed AI into every customer touchpoint",
                    "Build an AI-ready data platform",
                    "Upskill the entire workforce",
                    "Govern AI responsibly and transparently",
                ]
            ),
            speaker_notes="Cover the five pillars briefly — we'll deep-dive each.",
        ),
        # 3. Stats
        Slide(
            id=SlideID.generate(),
            type=SlideType.stats,
            heading="Baseline Metrics",
            content=SlideContent(
                stats=[
                    StatItem(value="23%", label="Automation Rate", desc="of processes automated"),
                    StatItem(value="$4.2M", label="Annual Savings", desc="from existing AI"),
                    StatItem(value="68%", label="Data Quality", desc="structured data coverage"),
                    StatItem(value="1,200", label="AI Workflows", desc="deployed in production"),
                ]
            ),
            speaker_notes="Establish where we are today before discussing where we're going.",
        ),
        # 4. Steps (Process)
        Slide(
            id=SlideID.generate(),
            type=SlideType.steps,
            heading="AI Adoption Framework",
            content=SlideContent(
                steps=[
                    "Assess current capabilities and data maturity",
                    "Identify high-ROI AI use cases",
                    "Build proof-of-concept models",
                    "Scale successful pilots to production",
                    "Monitor, iterate, and govern continuously",
                ]
            ),
            speaker_notes="Our five-phase approach ensures controlled, ROI-positive adoption.",
        ),
        # 5. Cards  ← will be EDITED
        Slide(
            id=SlideID.generate(),
            type=SlideType.cards,
            heading="Core Capabilities",
            content=SlideContent(
                cards=[
                    CardItem(icon="🧠", title="GenAI Platform", desc="Enterprise LLM deployment"),
                    CardItem(icon="📊", title="Data Fabric", desc="Unified data infrastructure"),
                    CardItem(icon="🤖", title="ML Ops", desc="Model lifecycle management"),
                    CardItem(icon="🔒", title="AI Governance", desc="Responsible AI framework"),
                ]
            ),
            speaker_notes="These four capabilities form the backbone of our AI platform.",
        ),
        # 6. Timeline
        Slide(
            id=SlideID.generate(),
            type=SlideType.timeline,
            heading="Delivery Roadmap",
            content=SlideContent(
                timeline=[
                    TimelineItem(date="Q1 2026", title="Foundation", desc="Data platform & governance"),
                    TimelineItem(date="Q2 2026", title="Pilot Phase", desc="10 GenAI use cases"),
                    TimelineItem(date="Q3 2026", title="Scale-Up", desc="Enterprise rollout"),
                    TimelineItem(date="Q4 2026", title="Optimize", desc="ROI measurement & tuning"),
                ]
            ),
            speaker_notes="Four quarters from foundation to optimisation.",
        ),
        # 7. Compare
        Slide(
            id=SlideID.generate(),
            type=SlideType.compare,
            heading="Build vs Buy",
            content=SlideContent(
                left_column=CompareColumn(
                    heading="Build In-House",
                    points=[
                        "Full IP ownership",
                        "Custom to our workflows",
                        "Long time to value",
                        "High talent requirement",
                    ],
                ),
                right_column=CompareColumn(
                    heading="Buy / Partner",
                    points=[
                        "Rapid deployment",
                        "Vendor lock-in risk",
                        "Lower upfront cost",
                        "Proven technology",
                    ],
                ),
            ),
            speaker_notes="Hybrid approach recommended — buy platform, build use-cases.",
        ),
        # 8. Chart
        Slide(
            id=SlideID.generate(),
            type=SlideType.chart,
            heading="Projected ROI Impact",
            content=SlideContent(
                chart_type=ChartType.column,
                chart_categories=["2024 Baseline", "2025 Est.", "2026 Target"],
                chart_series=[
                    ChartSeries(name="Cost Savings ($M)", values=[4.2, 8.5, 18.0]),
                    ChartSeries(name="Revenue Uplift ($M)", values=[0.0, 3.0, 12.0]),
                ],
                chart_title="AI-Driven Financial Impact",
            ),
            speaker_notes="ROI inflects sharply in 2026 as scaled platforms mature.",
        ),
    ]

    meta = DeckMeta(
        title="AI Strategy 2026",
        description="Enterprise AI transformation strategy and roadmap",
        author="Demo User",
        created_at=datetime.utcnow(),
        updated_at=datetime.utcnow(),
        version=1,
        change_log=[
            ChangeEntry(
                version=1,
                author="Demo User",
                action="create",
                slide_ids_affected=[s.id for s in slides],
                description="Initial deck creation",
            )
        ],
    )

    theme = DeckTheme(
        colors=ThemeColors(
            accent="#6366F1",
            accent2="#A78BFA",
        ),
        fonts=ThemeFonts(
            heading="Calibri Light",
            body="Calibri",
        ),
    )

    return Deck(
        meta=meta,
        theme=theme,
        plan=DeckPlan(slides=slides),
    )


# ── Main demo flow ─────────────────────────────────────────────────────────────

def main():
    print("\n" + "═" * 72)
    print("  Barq PPTX Engine — End-to-End Demo")
    print("═" * 72)

    # ── Step 1: Build components ───────────────────────────────────────────────
    _sep("Step 1: Initialize Engine")
    registry = SlideTypeRegistry.build_default()
    renderer = PresentationRenderer(registry)
    edit_engine = EditEngine(renderer)
    validator = DeckValidator()

    # Use a temp file for the demo store
    db_fd, db_path = tempfile.mkstemp(suffix=".db", prefix="barq_demo_")
    os.close(db_fd)
    store = DeckStore(db_path=db_path)
    print(f"  Registry: {len(list(registry._renderers))} renderers registered")
    print(f"  Store:    {db_path}")

    # ── Step 2: Create the deck ────────────────────────────────────────────────
    _sep("Step 2: Create 'AI Strategy 2026' Deck (8 slides)")
    deck = _build_ai_strategy_deck()
    _print_deck_summary(deck)

    # Render
    print("\n  Rendering...")
    pptx_bytes = renderer.render_to_bytes(deck)
    print(f"  Rendered: {len(pptx_bytes):,} bytes")

    # Store
    store.create(deck, pptx_bytes)
    print("  Stored in DeckStore (v1)")

    # ── Step 3: Validate ───────────────────────────────────────────────────────
    _sep("Step 3: Validate")
    validation = validator.validate(deck, pptx_bytes)
    status = "PASSED" if validation.passed else "FAILED"
    print(f"  Status   : {status}")
    print(f"  Errors   : {validation.error_count}")
    print(f"  Warnings : {validation.warning_count}")
    print(f"  Info     : {validation.info_count}")
    if validation.issues:
        for issue in validation.issues[:5]:
            print(f"    [{issue.severity.value.upper():7s}] {issue.code}: {issue.message[:60]}")

    # ── Step 4: Edit slides 2 and 5 ───────────────────────────────────────────
    _sep("Step 4: Edit Slides 2 and 5")

    slide_2 = deck.plan.slides[1]  # "Strategic Pillars" → bullets
    slide_5 = deck.plan.slides[4]  # "Core Capabilities" → cards

    updated_slide_2 = slide_2.model_copy(update={
        "heading": "2026 Strategic Priorities",
        "content": SlideContent(
            points=[
                "Deploy GenAI across all 12 business units by Q3",
                "Reduce manual operations cost by 35% using automation",
                "Launch AI-powered customer support assistant",
                "Establish Centre of Excellence with 50 AI engineers",
                "Achieve ISO 42001 AI Management certification",
                "Generate $12M incremental revenue through AI products",
            ]
        ),
        "speaker_notes": "Updated with board-approved 2026 OKRs.",
    })

    updated_slide_5 = slide_5.model_copy(update={
        "heading": "Platform Capabilities — Updated",
        "content": SlideContent(
            cards=[
                CardItem(icon="⚡", title="GenAI Studio", desc="Visual AI workflow builder"),
                CardItem(icon="🔗", title="Integration Hub", desc="Connect 200+ enterprise systems"),
                CardItem(icon="📈", title="Insight Engine", desc="Real-time AI-powered analytics"),
                CardItem(icon="🛡️", title="Trust & Safety", desc="Bias detection & audit trails"),
                CardItem(icon="🌐", title="Multi-Cloud", desc="AWS, Azure, GCP deployment"),
                CardItem(icon="🎓", title="AI Academy", desc="Workforce upskilling platform"),
            ]
        ),
        "speaker_notes": "Updated to reflect expanded platform after Q1 2026 acquisition.",
    })

    edit_req = EditRequest(
        deck_id=deck.meta.id,
        slide_ids=[slide_2.id, slide_5.id],
        updated_slides=[updated_slide_2, updated_slide_5],
        author="Demo User",
        description="Updated strategic priorities and platform capabilities for board review",
        confirm=False,
    )

    print(f"  Targeting slides: {slide_2.id}, {slide_5.id}")
    edit_result = edit_engine.apply_edit(deck, edit_req, pptx_bytes)
    print(f"  Changed: {edit_result.changed_slide_ids}")
    print(f"  New version: {edit_result.version}")

    # Persist
    updated_deck = edit_result.deck
    updated_pptx = edit_result.pptx_bytes
    change = updated_deck.meta.change_log[-1]
    store.update(deck.meta.id, updated_deck, updated_pptx, change,
                 require_confirmation=False)
    print("  Stored in DeckStore (v2)")

    # ── Step 5: Verify no-touch preservation ──────────────────────────────────
    _sep("Step 5: Verify No-Touch Preservation")

    orig_prs = Presentation(io.BytesIO(pptx_bytes))
    result_prs = Presentation(io.BytesIO(updated_pptx))

    touched = {slide_2.id, slide_5.id}
    preserved_count = 0
    modified_count = 0

    for idx, (orig_slide, res_slide) in enumerate(zip(orig_prs.slides, result_prs.slides)):
        from pptx_engine.edit_engine import _read_slide_id
        sid = _read_slide_id(orig_slide)
        if sid in touched:
            print(f"  [EDITED]      slide[{idx+1}] {sid}")
            modified_count += 1
        else:
            print(f"  [PRESERVED]   slide[{idx+1}] {sid}")
            preserved_count += 1

    print(f"\n  Summary: {modified_count} edited, {preserved_count} preserved")

    # ── Step 6: Export final PPTX ─────────────────────────────────────────────
    _sep("Step 6: Export Final PPTX")

    output_path = "/tmp/demo_deck.pptx"
    with open(output_path, "wb") as f:
        f.write(updated_pptx)

    file_size = os.path.getsize(output_path)
    final_prs = Presentation(output_path)
    print(f"  Output   : {output_path}")
    print(f"  Size     : {file_size:,} bytes")
    print(f"  Slides   : {len(final_prs.slides)}")
    print(f"  Version  : v{updated_deck.meta.version}")

    # ── Step 7: Print change log ───────────────────────────────────────────────
    _sep("Step 7: Change Log")
    change_log = store.get_change_log(deck.meta.id)
    _print_change_log(change_log)

    # ── Summary ───────────────────────────────────────────────────────────────
    _sep("Demo Complete")
    print(f"  PPTX exported to: {output_path}")
    print(f"  {len(final_prs.slides)} slides | v{updated_deck.meta.version} | {file_size:,} bytes")
    print()

    # Cleanup temp DB
    try:
        os.unlink(db_path)
    except Exception:
        pass


if __name__ == "__main__":
    main()
