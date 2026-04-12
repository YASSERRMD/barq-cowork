#!/usr/bin/env python3
"""Convert PPTX/DOCX files to HTML preview.

Usage: python3 preview_doc.py <file_path>
Outputs HTML to stdout.
"""
import sys
import os

def preview_pptx(path: str) -> str:
    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation(path)

    # Extract theme colors from first slide background
    bg_color = "#0F172A"
    accent_color = "#6366F1"
    surface_color = "#1E293B"
    text_color = "#F8FAFC"
    muted_color = "#94A3B8"

    # Try to detect colors from the actual slides
    for slide in prs.slides:
        bg = slide.background
        if bg and bg.fill and bg.fill.type is not None:
            try:
                c = bg.fill.fore_color.rgb
                bg_color = f"#{c}"
                # Derive surface as slightly lighter
                r, g, b = int(str(c)[:2], 16), int(str(c)[2:4], 16), int(str(c)[4:6], 16)
                sr, sg, sb = min(r + 20, 255), min(g + 20, 255), min(b + 20, 255)
                surface_color = f"#{sr:02x}{sg:02x}{sb:02x}"
            except Exception:
                pass
        # Try to find accent from shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color and run.font.color.rgb:
                            c = str(run.font.color.rgb)
                            r, g, b = int(c[:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                            # If it's a bright/saturated color, use as accent
                            if max(r, g, b) > 100 and (max(r, g, b) - min(r, g, b)) > 50:
                                accent_color = f"#{c}"
                                break
        break  # only check first slide

    slides_html = []
    for i, slide in enumerate(prs.slides):
        shapes_html = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # Detect heading vs body by font size
                    font_size = None
                    font_bold = False
                    font_color = text_color
                    for run in para.runs:
                        if run.font.size:
                            font_size = run.font.size.pt
                        if run.font.bold:
                            font_bold = True
                        if run.font.color and run.font.color.rgb:
                            font_color = f"#{run.font.color.rgb}"

                    if font_size and font_size >= 24 or font_bold:
                        shapes_html.append(
                            f'<h2 style="color:{font_color};margin:0 0 8px;font-size:{min(font_size or 28, 32)}px;'
                            f'font-weight:700;letter-spacing:-0.02em">{_esc(text)}</h2>'
                        )
                    elif font_size and font_size >= 18:
                        shapes_html.append(
                            f'<h3 style="color:{font_color};margin:0 0 6px;font-size:{min(font_size, 22)}px;'
                            f'font-weight:600">{_esc(text)}</h3>'
                        )
                    else:
                        shapes_html.append(
                            f'<p style="color:{muted_color};margin:0 0 4px;font-size:14px;'
                            f'line-height:1.6">{_esc(text)}</p>'
                        )

            if shape.has_table:
                tbl = shape.table
                rows_html = []
                for row_idx, row in enumerate(tbl.rows):
                    cells = []
                    tag = "th" if row_idx == 0 else "td"
                    for cell in row.cells:
                        style = (f'padding:8px 12px;border:1px solid {surface_color};'
                                 f'font-size:13px;color:{text_color if row_idx > 0 else accent_color};'
                                 f'{"font-weight:600;" if row_idx == 0 else ""}'
                                 f'{"background:" + surface_color + ";" if row_idx == 0 else ""}')
                        cells.append(f'<{tag} style="{style}">{_esc(cell.text)}</{tag}>')
                    rows_html.append(f'<tr>{"".join(cells)}</tr>')
                shapes_html.append(
                    f'<table style="width:100%;border-collapse:collapse;margin:8px 0">'
                    f'{"".join(rows_html)}</table>'
                )

            if hasattr(shape, "chart"):
                shapes_html.append(
                    f'<div style="background:{surface_color};border:1px solid rgba(255,255,255,0.1);'
                    f'border-radius:8px;padding:16px;margin:8px 0;text-align:center">'
                    f'<span style="font-size:32px">📊</span>'
                    f'<p style="color:{muted_color};font-size:13px;margin:4px 0 0">Chart (view in PowerPoint)</p>'
                    f'</div>'
                )

        content = "\n".join(shapes_html) if shapes_html else (
            f'<p style="color:{muted_color};font-style:italic;font-size:13px">Empty slide</p>'
        )

        slides_html.append(
            f'<div style="background:{bg_color};border:1px solid rgba(255,255,255,0.08);'
            f'border-radius:12px;padding:28px 32px;margin-bottom:16px;'
            f'box-shadow:0 2px 12px rgba(0,0,0,0.3)">'
            f'<div style="display:flex;align-items:center;gap:8px;margin-bottom:16px">'
            f'<span style="background:{accent_color};color:#fff;font-size:11px;font-weight:700;'
            f'border-radius:6px;padding:2px 8px;min-width:24px;text-align:center">{i+1}</span>'
            f'<span style="font-size:11px;color:{muted_color}">Slide {i+1} of {len(prs.slides)}</span>'
            f'</div>'
            f'{content}'
            f'</div>'
        )

    return (
        f'<!DOCTYPE html><html><head><meta charset="utf-8">'
        f'<style>*{{box-sizing:border-box;margin:0;padding:0}}'
        f'body{{background:{bg_color};font-family:Inter,system-ui,sans-serif;'
        f'padding:24px;color:{text_color}}}</style></head><body>'
        f'{"".join(slides_html)}'
        f'</body></html>'
    )


def preview_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)

    bg = "#0F172A"
    surface = "#1E293B"
    accent = "#6366F1"
    text = "#F8FAFC"
    muted = "#94A3B8"

    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        style_name = (para.style.name or "").lower()
        if "heading 1" in style_name or "title" in style_name:
            parts.append(f'<h1 style="color:{accent};font-size:28px;font-weight:700;'
                         f'margin:24px 0 12px;letter-spacing:-0.02em;'
                         f'border-bottom:2px solid {surface};padding-bottom:8px">{_esc(t)}</h1>')
        elif "heading 2" in style_name:
            parts.append(f'<h2 style="color:{text};font-size:22px;font-weight:600;'
                         f'margin:20px 0 8px">{_esc(t)}</h2>')
        elif "heading" in style_name:
            parts.append(f'<h3 style="color:{text};font-size:18px;font-weight:600;'
                         f'margin:16px 0 6px">{_esc(t)}</h3>')
        elif t.startswith("•") or t.startswith("-") or "list" in style_name:
            bullet = t.lstrip("•-").strip()
            parts.append(f'<div style="display:flex;gap:8px;margin:4px 0 4px 16px;font-size:14px;'
                         f'color:{muted};line-height:1.6">'
                         f'<span style="color:{accent};flex-shrink:0">•</span>{_esc(bullet)}</div>')
        else:
            parts.append(f'<p style="color:{muted};font-size:14px;line-height:1.7;'
                         f'margin:6px 0">{_esc(t)}</p>')

    for table in doc.tables:
        rows_html = []
        for row_idx, row in enumerate(table.rows):
            cells = []
            tag = "th" if row_idx == 0 else "td"
            for cell in row.cells:
                style = (f'padding:8px 12px;border:1px solid {surface};font-size:13px;'
                         f'color:{text if row_idx > 0 else accent};'
                         f'{"font-weight:600;background:" + surface + ";" if row_idx == 0 else ""}')
                cells.append(f'<{tag} style="{style}">{_esc(cell.text)}</{tag}>')
            rows_html.append(f'<tr>{"".join(cells)}</tr>')
        parts.append(f'<table style="width:100%;border-collapse:collapse;margin:16px 0">'
                     f'{"".join(rows_html)}</table>')

    return (
        f'<!DOCTYPE html><html><head><meta charset="utf-8">'
        f'<style>*{{box-sizing:border-box;margin:0;padding:0}}'
        f'body{{background:{bg};font-family:Inter,system-ui,sans-serif;'
        f'padding:32px;max-width:800px;margin:0 auto;color:{text}}}</style></head><body>'
        f'{"".join(parts)}'
        f'</body></html>'
    )


def _esc(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def main():
    if len(sys.argv) < 2:
        print("Usage: preview_doc.py <file_path>", file=sys.stderr)
        sys.exit(1)

    path = sys.argv[1]
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pptx":
        print(preview_pptx(path))
    elif ext == ".docx":
        print(preview_docx(path))
    else:
        print(f"<html><body><p>Preview not available for {ext} files</p></body></html>")


if __name__ == "__main__":
    main()
