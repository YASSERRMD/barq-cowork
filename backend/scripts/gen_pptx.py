#!/usr/bin/env python3
"""
Barq Cowork — Professional PPTX Generator (python-pptx 1.0.x)
Reads JSON from stdin, writes .pptx bytes to stdout.

Supported layouts per slide:
  bullets   — numbered glassmorphism cards
  stats     — big KPI numbers + real bar/doughnut charts
  steps     — chevron flow diagram
  cards     — icon feature grid
  chart     — full-slide chart (column / bar / line / pie / doughnut)
  timeline  — horizontal milestone timeline
  compare   — two-column comparison (pros vs cons, before vs after)
"""

import json, sys, re
from io import BytesIO
from copy import deepcopy

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from lxml import etree

# ── Slide canvas ──────────────────────────────────────────────────────────────
W, H = 9144000, 6858000      # 10 × 7.5 inches in EMU

# ── Dark-mode color palette ───────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x17, 0x2A)   # #0F172A very dark navy
CARD     = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B dark slate
CARD2    = RGBColor(0x0D, 0x1B, 0x2E)   # even darker
WHITE    = RGBColor(0xF8, 0xFA, 0xFC)   # near-white
OFFWHITE = RGBColor(0xE2, 0xE8, 0xF0)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
LINE     = RGBColor(0x2D, 0x3F, 0x55)

def _h(hex6):
    h = hex6.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _v(c): return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

def _lighter(c, amt=60):
    return RGBColor(min(255,c[0]+amt), min(255,c[1]+amt), min(255,c[2]+amt))

def _darker(c, amt=40):
    return RGBColor(max(0,c[0]-amt), max(0,c[1]-amt), max(0,c[2]-amt))

# ── Slide helpers ─────────────────────────────────────────────────────────────

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color

def nofill(shape): shape.fill.background()
def noline(shape): shape.line.fill.background()

# ── XML primitives ────────────────────────────────────────────────────────────

def _spPr(shape):
    sp = shape._element
    return sp.find(qn("p:spPr"))

def _rm_fills(spPr):
    for t in ["a:solidFill","a:gradFill","a:noFill","a:blipFill","a:pattFill"]:
        e = spPr.find(qn(t))
        if e is not None: spPr.remove(e)

def xml_solid_alpha(shape, color, alpha_pct):
    spPr = _spPr(shape); _rm_fills(spPr)
    sf = etree.SubElement(spPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr")); sc.set("val", _v(color))
    al = etree.SubElement(sc, qn("a:alpha")); al.set("val", str(int(alpha_pct*1000)))

def xml_grad(shape, stops, angle_deg=90):
    spPr = _spPr(shape); _rm_fills(spPr)
    gf = etree.SubElement(spPr, qn("a:gradFill"))
    gs_lst = etree.SubElement(gf, qn("a:gsLst"))
    for pos, color in stops:
        gs = etree.SubElement(gs_lst, qn("a:gs")); gs.set("pos", str(int(pos*1000)))
        sc = etree.SubElement(gs, qn("a:srgbClr")); sc.set("val", _v(color))
    lin = etree.SubElement(gf, qn("a:lin"))
    lin.set("ang", str(int(angle_deg*60000))); lin.set("scaled","0")

def xml_shadow(shape, blur_pt=8, dist_pt=5, dir_deg=315, alpha=50):
    spPr = _spPr(shape)
    el = spPr.find(qn("a:effectLst"))
    if el is None: el = etree.SubElement(spPr, qn("a:effectLst"))
    old = el.find(qn("a:outerShdw"))
    if old is not None: el.remove(old)
    s = etree.SubElement(el, qn("a:outerShdw"))
    s.set("blurRad", str(int(blur_pt*12700)))
    s.set("dist",    str(int(dist_pt*12700)))
    s.set("dir",     str(int(dir_deg*60000)))
    s.set("algn","ctr"); s.set("rotWithShape","0")
    sc = etree.SubElement(s, qn("a:srgbClr")); sc.set("val","000000")
    etree.SubElement(sc, qn("a:alpha")).set("val", str(int(alpha*1000)))

def xml_stroke_alpha(shape, color, alpha_pct, w_pt=1.0):
    spPr = _spPr(shape)
    ln = spPr.find(qn("a:ln"))
    if ln is None: ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(w_pt*12700)))
    for c in list(ln): ln.remove(c)
    sf = etree.SubElement(ln, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr")); sc.set("val", _v(color))
    etree.SubElement(sc, qn("a:alpha")).set("val", str(int(alpha_pct*1000)))

def xml_round_corners(shape, adj=25000):
    spPr = _spPr(shape)
    pg = spPr.find(qn("a:prstGeom"))
    if pg is None: return
    av = pg.find(qn("a:avLst"))
    if av is None: av = etree.SubElement(pg, qn("a:avLst"))
    for g in av.findall(qn("a:gd")): av.remove(g)
    gd = etree.SubElement(av, qn("a:gd"))
    gd.set("name","adj"); gd.set("fmla",f"val {adj}")

def xml_set_geom(shape, prst):
    """Change a shape's preset geometry (e.g. 'chevron', 'pentagon', 'rightArrow')."""
    spPr = _spPr(shape)
    pg = spPr.find(qn("a:prstGeom"))
    if pg is not None:
        pg.set("prst", prst)
        av = pg.find(qn("a:avLst"))
        if av is None: av = etree.SubElement(pg, qn("a:avLst"))
        for g in list(av): av.remove(g)

# ── Chart XML helpers ─────────────────────────────────────────────────────────

def _get_or_add(parent, tag):
    """Get or add a child element by tag."""
    el = parent.find(tag)
    if el is None:
        el = etree.SubElement(parent, tag)
    return el

def _set_solid_fill(spPr, hex_val):
    """Set solid fill on a spPr element, removing any existing fill."""
    _rm_fills_el(spPr)
    sf = etree.SubElement(spPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", hex_val)

def _set_no_border(spPr):
    """Remove border from a spPr element."""
    ln = _get_or_add(spPr, qn("a:ln"))
    for c in list(ln): ln.remove(c)
    etree.SubElement(ln, qn("a:noFill"))

def _chart_dark_bg(chart, bg_color=CARD, area_color=BG):
    """Apply dark background to chart area and plot area using raw lxml XML."""
    # chart._element IS the chartSpace element in python-pptx 1.x
    cs = chart._element
    # chartSpace spPr is in 'c:' namespace
    spPr = _get_or_add(cs, qn("c:spPr"))
    _set_solid_fill(spPr, _v(area_color))
    _set_no_border(spPr)

    # plotArea spPr
    chart_el = cs.find(qn("c:chart"))
    if chart_el is not None:
        plotArea = chart_el.find(qn("c:plotArea"))
        if plotArea is not None:
            spPr2 = _get_or_add(plotArea, qn("c:spPr"))
            _set_solid_fill(spPr2, _v(bg_color))
            _set_no_border(spPr2)

def _rm_fills_el(el):
    for t in ["a:solidFill","a:gradFill","a:noFill","a:blipFill","a:pattFill"]:
        e = el.find(qn(t))
        if e is not None: el.remove(e)

def _style_axis(axis, label_color=MUTED, gridline_color=LINE, hide_gridlines=False):
    try:
        axis.tick_labels.font.color.rgb = label_color
        axis.tick_labels.font.size = Pt(9)
    except: pass
    try:
        if hide_gridlines:
            axis.major_gridlines.format.line.fill.background()
        else:
            axis.major_gridlines.format.line.color.rgb = gridline_color
            axis.major_gridlines.format.line.width = Pt(0.5)
    except: pass
    try:
        axis.format.line.fill.background()
    except: pass

# ── Shape factory ─────────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, color=None):
    s = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
    if color: solid(s, color)
    else: nofill(s)
    noline(s); return s

def rrect(slide, x, y, w, h, color=None, adj=22000):
    s = slide.shapes.add_shape(5, Emu(x), Emu(y), Emu(w), Emu(h))
    if color: solid(s, color)
    else: nofill(s)
    noline(s); xml_round_corners(s, adj); return s

def ellipse(slide, x, y, w, h, color=None):
    s = slide.shapes.add_shape(9, Emu(x), Emu(y), Emu(w), Emu(h))
    if color: solid(s, color)
    else: nofill(s)
    noline(s); return s

def chevron(slide, x, y, w, h, color):
    s = rect(slide, x, y, w, h, color)
    xml_set_geom(s, "chevron"); return s

def pentagon(slide, x, y, w, h, color):
    s = rect(slide, x, y, w, h, color)
    xml_set_geom(s, "pentagon"); return s

def arrow_right(slide, x, y, w, h, color):
    s = rect(slide, x, y, w, h, color)
    xml_set_geom(s, "rightArrow"); return s

# ── Text helpers ──────────────────────────────────────────────────────────────

def txbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))

def set_para(shape, text, size_pt, color, bold=False,
             font="Calibri Light", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
             italic=False, wrap=True):
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    # clear extra paragraphs
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    p.text = ""; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size_pt)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color

def label(slide, x, y, w, h, text, size_pt, color, bold=False,
          font="Calibri Light", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    s = txbox(slide, x, y, w, h)
    set_para(s, text, size_pt, color, bold, font, align, anchor)
    return s

# ── Background & header (shared by all content slides) ───────────────────────

def paint_bg(slide, accent):
    rect(slide, 0, 0, W, H, BG)
    # Large decorative circle top-right
    c1 = ellipse(slide, 6900000, -800000, 3000000, 3000000)
    xml_solid_alpha(c1, accent, 4)
    xml_stroke_alpha(c1, accent, 14, 2.5)
    # Inner circle
    c2 = ellipse(slide, 7500000, -300000, 1900000, 1900000)
    xml_solid_alpha(c2, accent, 3)
    xml_stroke_alpha(c2, accent, 8, 1.5)
    # Small dot bottom-left
    d = ellipse(slide, 180000, 6100000, 380000, 380000)
    xml_solid_alpha(d, accent, 10)

def paint_header(slide, heading, accent):
    # Gradient top bar
    tb = rect(slide, 0, 0, W, 14000)
    xml_grad(tb, [(0, accent), (100, _darker(accent, 30))], angle_deg=0)
    # Title
    t = txbox(slide, 457200, 55000, 7700000, 490000)
    set_para(t, heading, 31, WHITE, bold=True, font="Calibri Light",
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # Accent dot
    d = ellipse(slide, 380000, 277000, 55000, 55000, accent)
    # Divider
    div = rect(slide, 457200, 590000, 8229600, 8000, LINE)
    return 660000   # content_top_y

# ── Shared card renderer ──────────────────────────────────────────────────────

def card_bg(slide, x, y, w, h, accent, adj=22000):
    c = rrect(slide, x, y, w, h, CARD, adj)
    xml_stroke_alpha(c, accent, 18, 1.0)
    xml_shadow(c, blur_pt=10, dist_pt=5, dir_deg=315, alpha=45)
    return c

# ══════════════════════════════════════════════════════════════════════════════
# LAYOUT: BULLETS
# ══════════════════════════════════════════════════════════════════════════════

def layout_bullets(slide, points, accent, top):
    if not points: return
    n = len(points)
    avail = H - top - 180000
    two_col = n >= 5

    if two_col:
        cols, col_w, col_xs = 2, 3962400, [457200, 4724400]
        rows = (n+1)//2
    else:
        cols, col_w, col_xs = 1, 8229600, [457200]
        rows = n

    gap = 90000
    ch = min((avail - (rows-1)*gap)//rows, 1450000)
    total = rows*ch + (rows-1)*gap
    sy = top + (avail-total)//2

    for i, pt in enumerate(points):
        col, row = i%cols, i//cols
        cx, cy = col_xs[col], sy + row*(ch+gap)
        cw = col_w

        card_bg(slide, cx, cy, cw, ch, accent)

        # Left accent strip (gradient)
        strip = rect(slide, cx, cy+ch//8, 11000, ch*3//4)
        xml_grad(strip, [(0, _lighter(accent,60)), (50, accent), (100, _darker(accent,20))], 90)

        # Number badge
        ns = min(int(ch*0.52), 440000)
        nx, ny = cx+190000, cy+(ch-ns)//2
        nb = ellipse(slide, nx, ny, ns, ns)
        xml_grad(nb, [(0, _lighter(accent,50)), (100, accent)], 135)
        xml_shadow(nb, 4, 2, 315, 45)
        label(slide, nx, ny, ns, ns, str(i+1),
              13 if ns>350000 else 11, WHITE, bold=True, font="Calibri")

        # Point text
        tx = nx+ns+160000
        tw = cw-(tx-cx)-120000
        t = txbox(slide, tx, cy+55000, tw, ch-110000)
        set_para(t, pt, 15 if ch>900000 else 12, OFFWHITE,
                 font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════════════════════════════════════════════════
# LAYOUT: STATS  (big KPIs + real chart)
# ══════════════════════════════════════════════════════════════════════════════

def layout_stats(slide, stats_list, points, accent, top):
    items = stats_list or []
    # Fallback: parse points into stats
    if not items and points:
        for p in points:
            m = re.match(r'^([^\s]+)\s+(.*)', p.strip())
            items.append({"value": m.group(1) if m else p,
                          "label": m.group(2) if m else "",
                          "desc":  ""})
    items = items[:4]
    if not items: return
    n = len(items)

    # ── Row 1: KPI cards ──────────────────────────────────────────────────────
    avail_w = 8229600
    gap = 200000
    cw = (avail_w - (n-1)*gap)//n
    ch = 2600000
    sy = top + 60000
    sx = 457200

    has_pct = [bool(re.search(r'\d', it.get("value",""))) for it in items]
    numeric_vals = []

    for i, it in enumerate(items):
        cx = sx + i*(cw+gap)
        c = card_bg(slide, cx, sy, cw, ch, accent)
        # Top gradient strip
        ts = rect(slide, cx, sy, cw, ch//9)
        xml_grad(ts, [(0,accent),(100,_lighter(accent,50))], 0)

        val = it.get("value","")
        # Giant value
        label(slide, cx+15000, sy+ch//9+ch//18, cw-30000, ch*5//14,
              val, 46, accent, bold=True, font="Calibri Light")

        # Progress bar for %
        pct_m = re.search(r'(\d+(?:\.\d+)?)\s*%', val)
        bar_y = sy + ch//9 + ch*5//14 + ch//14
        if pct_m:
            pct = float(pct_m.group(1))/100.0
            bg_bar = rrect(slide, cx+50000, bar_y, cw-100000, ch//14, LINE, adj=50000)
            bw = int((cw-100000)*min(pct,1.0))
            if bw > 5000:
                fb = rrect(slide, cx+50000, bar_y, bw, ch//14, accent, adj=50000)
                xml_grad(fb, [(0,accent),(100,_lighter(accent,60))], 0)
            lbl_y = bar_y + ch//14 + ch//20
        else:
            lbl_y = bar_y

        # Label & desc
        label(slide, cx+15000, lbl_y, cw-30000, ch//6,
              it.get("label",""), 14, WHITE, bold=True, font="Calibri")
        if it.get("desc",""):
            label(slide, cx+15000, lbl_y+ch//6, cw-30000, ch//8,
                  it["desc"], 11, MUTED, font="Calibri")

        # Collect numeric value for chart
        nm = re.search(r'[\d,.]+', val)
        numeric_vals.append(float(nm.group().replace(",","")) if nm else 0)

    # ── Row 2: Bar chart showing all stats ────────────────────────────────────
    chart_top = sy + ch + 200000
    chart_h = H - chart_top - 200000
    if chart_h < 600000: return

    cd = CategoryChartData()
    cd.categories = [it.get("label","") for it in items]
    cd.add_series("Value", numeric_vals)

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(457200), Emu(chart_top), Emu(8229600), Emu(chart_h), cd)
    ch_obj = gf.chart
    ch_obj.has_legend = False
    _chart_dark_bg(ch_obj, CARD, BG)

    # Series color
    try:
        ser = ch_obj.series[0]
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = accent
        ser.data_labels.show_value = True
        ser.data_labels.font.color.rgb = WHITE
        ser.data_labels.font.size = Pt(10)
        ser.data_labels.font.bold = True
    except: pass

    try:
        _style_axis(ch_obj.category_axis, MUTED, LINE)
        _style_axis(ch_obj.value_axis,    MUTED, LINE)
    except: pass

# ══════════════════════════════════════════════════════════════════════════════
# LAYOUT: STEPS  (chevron flow)
# ══════════════════════════════════════════════════════════════════════════════

def layout_steps(slide, points, accent, top):
    if not points: return
    n = min(len(points), 6)
    two_rows = n >= 5
    rows_data = [points[:3], points[3:n]] if two_rows else [points[:n]]

    avail_h = H - top - 200000
    row_h = avail_h // len(rows_data)
    step_idx = 0

    for row_i, row_pts in enumerate(rows_data):
        rn = len(row_pts)
        # Each chevron + overlap
        chevron_w = (8229600 - 60000*(rn-1)) // rn
        chevron_h = int(row_h * 0.48)
        cy = top + row_i * row_h + (row_h - chevron_h)//2 - 40000
        text_y = cy + chevron_h + 80000
        text_h = row_h - chevron_h - 100000

        for i, pt in enumerate(row_pts):
            step_idx += 1
            cx = 457200 + i*(chevron_w + 60000)
            # Use chevron for middle/last, pentagon for first
            geom = "home" if i == 0 else "chevron"
            shp = rect(slide, cx, cy, chevron_w, chevron_h)
            # Gradient fill
            xml_grad(shp, [(0, _lighter(accent,30)), (100, _darker(accent,20))], 135)
            xml_set_geom(shp, geom)
            xml_shadow(shp, 8, 4, 315, 40)

            # Step number circle (centered top of chevron)
            ns = min(chevron_h//3, 380000)
            nx = cx + (chevron_w - ns)//2
            ny = cy - ns//2
            nb = ellipse(slide, nx, ny, ns, ns)
            xml_grad(nb, [(0, WHITE), (100, OFFWHITE)], 135)
            xml_shadow(nb, 3, 2, 315, 30)
            label(slide, nx, ny, ns, ns, str(step_idx),
                  12, accent, bold=True, font="Calibri")

            # Step text inside chevron
            inner_w = int(chevron_w * (0.7 if i < rn-1 else 0.85))
            t = txbox(slide, cx+20000, cy+20000, inner_w, chevron_h-40000)
            short = pt[:45]+"…" if len(pt)>45 else pt
            set_para(t, short, 12, WHITE, font="Calibri",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            # Full description below
            if len(pt) > 10:
                t2 = txbox(slide, cx, text_y, chevron_w, text_h)
                set_para(t2, pt, 11, MUTED, font="Calibri",
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# ══════════════════════════════════════════════════════════════════════════════
# LAYOUT: CARDS  (icon feature grid)
# ══════════════════════════════════════════════════════════════════════════════

# Icon glyphs in Segoe UI Symbol — safe cross-platform
_ICONS = ["●", "◆", "▲", "★", "◉", "⬡", "◈", "▣"]

def layout_cards(slide, points, accent, top):
    if not points: return
    n = min(len(points), 6)
    cols = 3 if n >= 5 else 2
    rows = (n+cols-1)//cols

    avail_h = H - top - 180000
    gx, gy = 200000, 160000
    cw = (8229600 - (cols-1)*gx) // cols
    ch = min((avail_h - (rows-1)*gy)//rows, 2200000)
    total_h = rows*ch + (rows-1)*gy
    sy = top + (avail_h-total_h)//2

    for i, pt in enumerate(points[:n]):
        col, row = i%cols, i//cols
        cx = 457200 + col*(cw+gx)
        cy = sy + row*(ch+gy)

        c = card_bg(slide, cx, cy, cw, ch, accent, adj=22000)

        # Gradient top strip
        ts = rrect(slide, cx, cy, cw, ch//8, adj=0)
        xml_grad(ts, [(0,accent),(60,_lighter(accent,60)),(100,accent)], 0)
        # Round top corners of strip manually — just use rect
        ts2 = rect(slide, cx, cy+ch//16, cw, ch//16, CARD)  # cover bottom half of strip

        # Icon circle
        icon_d = min(cw//4, 560000)
        ix = cx + (cw-icon_d)//2
        iy = cy + ch//8 + 70000
        ic = ellipse(slide, ix, iy, icon_d, icon_d)
        xml_solid_alpha(ic, accent, 18)
        xml_stroke_alpha(ic, accent, 45, 1.5)

        # Icon character
        label(slide, ix, iy, icon_d, icon_d,
              _ICONS[i % len(_ICONS)], 20, accent, bold=True,
              font="Segoe UI Symbol")

        # Title
        parts = pt.split(":",1)
        title_s = parts[0].strip()
        body_s  = parts[1].strip() if len(parts)>1 else ""

        ty = iy + icon_d + 60000
        label(slide, cx+30000, ty, cw-60000, ch//5,
              title_s, 14, WHITE, bold=True, font="Calibri",
              align=PP_ALIGN.CENTER)

        if body_s:
            by = ty + ch//5 + 20000
            rem = cy + ch - by - 60000
            if rem > 50000:
                label(slide, cx+35000, by, cw-70000, rem,
                      body_s, 11, MUTED, font="Calibri",
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# ══════════════════════════════════════════════════════════════════════════════
# LAYOUT: CHART  (dedicated full-slide chart)
# ══════════════════════════════════════════════════════════════════════════════

_CHART_TYPES = {
    "column":   XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar":      XL_CHART_TYPE.BAR_CLUSTERED,
    "line":     XL_CHART_TYPE.LINE_MARKERS,
    "pie":      XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area":     XL_CHART_TYPE.AREA,
}

_SERIES_COLORS = [
    RGBColor(0x63,0x66,0xF1),  # indigo
    RGBColor(0x06,0xB6,0xD4),  # cyan
    RGBColor(0x10,0xB9,0x81),  # green
    RGBColor(0xF5,0x9E,0x0B),  # amber
    RGBColor(0xF4,0x3F,0x5E),  # rose
    RGBColor(0x8B,0x5C,0xF6),  # purple
]

def layout_chart(slide, slide_data, accent, top):
    chart_type_key = slide_data.get("chart_type","column").lower()
    ct = _CHART_TYPES.get(chart_type_key, XL_CHART_TYPE.COLUMN_CLUSTERED)
    cats = slide_data.get("categories", [])
    series_data = slide_data.get("series", [])

    # Fall back: build from points
    if not cats and not series_data:
        pts = slide_data.get("points", [])
        for p in pts:
            m = re.match(r'^(.*?):\s*([\d,.]+)', p)
            if m:
                cats.append(m.group(1).strip())
                series_data = [{"name":"Value","values": [float(v.replace(",",""))
                    for v in [re.search(r'[\d,.]+',q).group() for q in pts
                               if re.search(r'[\d,.]+',q)]]}]
                break
        if not cats:
            cats = [f"Item {i+1}" for i in range(len(pts))]
            series_data = [{"name":"Value","values":[i+1 for i in range(len(pts))]}]

    cd = CategoryChartData()
    cd.categories = cats
    for s in series_data:
        cd.add_series(s.get("name",""), s.get("values",[]))

    ch_h = H - top - 150000
    gf = slide.shapes.add_chart(
        ct, Emu(457200), Emu(top), Emu(8229600), Emu(ch_h), cd)
    ch_obj = gf.chart
    ch_obj.has_legend = len(series_data) > 1
    if ch_obj.has_legend:
        try:
            ch_obj.legend.position = XL_LEGEND_POSITION.BOTTOM
            ch_obj.legend.include_in_layout = False
            ch_obj.legend.font.color.rgb = MUTED
            ch_obj.legend.font.size = Pt(10)
        except: pass

    _chart_dark_bg(ch_obj, CARD, BG)

    # Color each series
    for idx, ser in enumerate(ch_obj.series):
        col = _SERIES_COLORS[idx % len(_SERIES_COLORS)]
        try:
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = col
        except: pass
        try:
            if ct not in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT):
                ser.data_labels.show_value = True
                ser.data_labels.font.color.rgb = WHITE
                ser.data_labels.font.size = Pt(9)
        except: pass

    # Style line marker
    if ct == XL_CHART_TYPE.LINE_MARKERS:
        try:
            for ser in ch_obj.series:
                ser.smooth = True
        except: pass

    # Axes
    try:
        _style_axis(ch_obj.category_axis, MUTED, LINE)
        _style_axis(ch_obj.value_axis,    MUTED, LINE)
    except: pass

    # Pie/doughnut: color slices individually
    if ct in (XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT):
        try:
            for idx, pt in enumerate(ch_obj.series[0].points):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = _SERIES_COLORS[idx % len(_SERIES_COLORS)]
            ch_obj.series[0].data_labels.show_value = True
            ch_obj.series[0].data_labels.show_category_name = True
            ch_obj.series[0].data_labels.font.color.rgb = WHITE
            ch_obj.series[0].data_labels.font.size = Pt(10)
        except: pass
        if ct == XL_CHART_TYPE.DOUGHNUT:
            try:
                ch_obj.plots[0].hole_size = 65
            except: pass

# ══════════════════════════════════════════════════════════════════════════════
# LAYOUT: TIMELINE
# ══════════════════════════════════════════════════════════════════════════════

def layout_timeline(slide, points, accent, top):
    if not points: return
    n = min(len(points), 7)
    pts = points[:n]

    avail_h = H - top - 200000
    # Center line Y
    line_y = top + avail_h//2 - 20000

    # Background: horizontal accent line
    spine = rect(slide, 457200, line_y+20000, 8229600, 10000, LINE)

    # Distribute milestones
    spacing = 8229600 // (n if n>1 else 1)
    for i, pt in enumerate(pts):
        mx = 457200 + i*spacing + spacing//2
        above = (i % 2 == 0)   # alternate above/below

        # Milestone dot (accent circle)
        ds = 220000
        dx = mx - ds//2
        dy = line_y + 20000 - ds//2
        d = ellipse(slide, dx, dy, ds, ds, accent)
        xml_shadow(d, 5, 3, 315, 45)

        # Connector line to dot
        cn_h = avail_h//3 - 80000
        if above:
            conn = rect(slide, mx-5000, line_y+20000-cn_h, 10000, cn_h, LINE)
            # Date / label above
            pdate, pdesc = _split_timeline_pt(pt)
            label(slide, mx-500000, line_y-cn_h-200000, 1000000, 180000,
                  pdate, 13, accent, bold=True, font="Calibri", align=PP_ALIGN.CENTER)
            t = txbox(slide, mx-500000, line_y-cn_h-380000, 1000000, 200000)
            set_para(t, pdesc, 10, OFFWHITE, font="Calibri",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        else:
            conn = rect(slide, mx-5000, line_y+ds+40000, 10000, cn_h-ds, LINE)
            pdate, pdesc = _split_timeline_pt(pt)
            label(slide, mx-500000, line_y+ds+cn_h+80000, 1000000, 180000,
                  pdate, 13, accent, bold=True, font="Calibri", align=PP_ALIGN.CENTER)
            t = txbox(slide, mx-500000, line_y+ds+cn_h+260000, 1000000, 200000)
            set_para(t, pdesc, 10, OFFWHITE, font="Calibri",
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

def _split_timeline_pt(pt):
    m = re.match(r'^(\d{4}(?:[–\-]\d{4})?|[A-Z][a-z]+\s+\d{4})[:\-–]\s*(.*)', pt)
    if m: return m.group(1), m.group(2)
    parts = pt.split(" ",1)
    return parts[0], parts[1] if len(parts)>1 else ""

# ══════════════════════════════════════════════════════════════════════════════
# LAYOUT: COMPARE  (two-column side-by-side)
# ══════════════════════════════════════════════════════════════════════════════

def layout_compare(slide, slide_data, accent, top):
    left_pts  = slide_data.get("left",  [])
    right_pts = slide_data.get("right", [])
    left_title  = slide_data.get("left_title",  "Before")
    right_title = slide_data.get("right_title", "After")

    # Fall back: split points evenly
    if not left_pts and not right_pts:
        pts = slide_data.get("points",[])
        mid = len(pts)//2
        left_pts, right_pts = pts[:mid], pts[mid:]

    avail_h = H - top - 100000
    col_w = (8229600 - 200000)//2
    col_h = avail_h

    for col_i, (title, items, col_accent) in enumerate([
        (left_title,  left_pts,  _darker(accent, 20)),
        (right_title, right_pts, accent),
    ]):
        cx = 457200 + col_i*(col_w+200000)

        # Column card
        c = rrect(slide, cx, top, col_w, col_h, CARD, adj=18000)
        xml_stroke_alpha(c, col_accent, 30, 1.5)
        xml_shadow(c, 12, 6, 315, 40)

        # Header
        hdr = rrect(slide, cx, top, col_w, 340000, col_accent, adj=0)
        label(slide, cx+40000, top+60000, col_w-80000, 220000,
              title, 18, WHITE, bold=True, font="Calibri Light",
              align=PP_ALIGN.CENTER)

        # Items
        item_h = 260000
        item_gap = 50000
        for j, it in enumerate(items[:8]):
            iy = top + 380000 + j*(item_h+item_gap)
            if iy + item_h > top + col_h - 60000: break

            # Check/cross mark
            mark = "✓" if col_i == 1 else "✕"
            mark_c = RGBColor(0x10,0xB9,0x81) if col_i==1 else RGBColor(0xF4,0x3F,0x5E)
            mb = ellipse(slide, cx+40000, iy+(item_h-180000)//2, 180000, 180000)
            xml_solid_alpha(mb, mark_c, 18)
            label(slide, cx+40000, iy+(item_h-180000)//2, 180000, 180000,
                  mark, 11, mark_c, bold=True, font="Segoe UI Symbol")

            t = txbox(slide, cx+260000, iy, col_w-320000, item_h)
            set_para(t, it, 13, OFFWHITE, font="Calibri",
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════════════════════════════════════════════════
# COVER SLIDE
# ══════════════════════════════════════════════════════════════════════════════

def build_cover(prs, title, subtitle, accent):
    slide = blank_slide(prs)
    lighter = _lighter(accent, 80)
    darker  = _darker(accent, 40)

    rect(slide, 0, 0, W, H, BG)

    # Big decorative circles top-right
    c1 = ellipse(slide, 5700000, -1300000, 4400000, 4400000)
    xml_solid_alpha(c1, accent, 4)
    xml_stroke_alpha(c1, accent, 10, 2.5)
    c2 = ellipse(slide, 6400000, -700000, 3000000, 3000000)
    xml_solid_alpha(c2, lighter, 3)
    xml_stroke_alpha(c2, lighter, 7, 1.5)
    c3 = ellipse(slide, 7200000, -200000, 1800000, 1800000)
    xml_solid_alpha(c3, accent, 5)

    # Small dot cluster bottom-left
    for ox, oy, od in [(150000,6200000,350000),(560000,6420000,180000),(820000,6280000,100000)]:
        d = ellipse(slide, ox, oy, od, od)
        xml_solid_alpha(d, accent, 8)

    # Left vertical accent bar (gradient)
    bar = rect(slide, 420000, 1000000, 18000, 4600000)
    xml_grad(bar, [(0,lighter),(50,accent),(100,darker)], 90)

    # Horizontal rule
    rule = rect(slide, 600000, 3800000, 3500000, 7000, LINE)

    # Title text
    t = txbox(slide, 660000, 1050000, 7700000, 2600000)
    set_para(t, title, 44, WHITE, bold=True, font="Calibri Light",
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # Subtitle
    s = txbox(slide, 660000, 3750000, 5800000, 600000)
    set_para(s, subtitle, 20, MUTED, font="Calibri",
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom gradient bar
    bb = rect(slide, 0, H-14000, W, 14000)
    xml_grad(bb, [(0,darker),(50,accent),(100,darker)], 0)

# ══════════════════════════════════════════════════════════════════════════════
# CONTENT SLIDE DISPATCHER
# ══════════════════════════════════════════════════════════════════════════════

def _auto_layout(heading, points, stats):
    if stats: return "stats"
    h = heading.lower()
    for kw in ["step","process","how ","workflow","pipeline","roadmap","phase","stage","journey"]:
        if kw in h and len(points) >= 2: return "steps"
    for kw in ["timeline","history","evolution","milestones","years","decade"]:
        if kw in h: return "timeline"
    for kw in ["vs","compare","versus","pros","cons","before","after","difference"]:
        if kw in h: return "compare"
    num_re = re.compile(r'^\s*[\d$€£>~]')
    if sum(1 for p in points if num_re.match(p)) >= max(1,(len(points)+1)//2):
        return "stats"
    if 4 <= len(points) <= 6 and all(len(p)<90 for p in points): return "cards"
    return "bullets"

def build_content_slide(prs, sd, accent):
    slide = blank_slide(prs)
    paint_bg(slide, accent)
    top = paint_header(slide, sd.get("heading",""), accent)

    layout = sd.get("layout","") or _auto_layout(
        sd.get("heading",""), sd.get("points",[]), sd.get("stats",[]))

    if   layout == "stats":    layout_stats(slide, sd.get("stats",[]), sd.get("points",[]), accent, top)
    elif layout == "steps":    layout_steps(slide, sd.get("points",[]), accent, top)
    elif layout == "cards":    layout_cards(slide, sd.get("points",[]), accent, top)
    elif layout == "chart":    layout_chart(slide, sd, accent, top)
    elif layout == "timeline": layout_timeline(slide, sd.get("points",[]), accent, top)
    elif layout == "compare":  layout_compare(slide, sd, accent, top)
    else:                      layout_bullets(slide, sd.get("points",[]), accent, top)

# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    try:
        data = json.loads(sys.stdin.read())
    except Exception as e:
        sys.stderr.write(f"JSON parse error: {e}\n"); sys.exit(1)

    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)

    accent = _h(data.get("accent","6366F1"))
    build_cover(prs, data.get("title",""), data.get("subtitle",""), accent)
    for sd in data.get("slides",[]):
        build_content_slide(prs, sd, accent)

    buf = BytesIO()
    prs.save(buf)
    sys.stdout.buffer.write(buf.getvalue())

if __name__ == "__main__":
    main()
